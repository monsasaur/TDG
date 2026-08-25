from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette — minimal monochrome + single accent ──────────
BLACK  = RGBColor(0x0D, 0x0D, 0x0D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x2E, 0x86, 0xDE)
LGRAY  = RGBColor(0xF7, 0xF8, 0xFA)
MGRAY  = RGBColor(0xE2, 0xE6, 0xED)
DGRAY  = RGBColor(0x6B, 0x74, 0x80)
RED    = RGBColor(0xE6, 0x39, 0x46)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=WHITE, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h,
        size=14, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size  = Pt(size)
    r.font.bold  = bold
    r.font.color.rgb = color
    r.font.italic = italic
    r.font.name  = "Helvetica Neue"
    return tb

def para(tf, text, size=13, bold=False, color=DGRAY,
         align=PP_ALIGN.LEFT, space=8, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.italic = italic
    r.font.name = "Helvetica Neue"

def bg(slide): rect(slide, 0, 0, W, H, WHITE)

def header(slide, label, title):
    """Minimal header: small label + large title + thin rule."""
    txt(slide, label, Inches(0.7), Inches(0.5), Inches(10), Inches(0.35),
        size=10, bold=True, color=BLUE)
    txt(slide, title, Inches(0.7), Inches(0.8), Inches(11.5), Inches(0.75),
        size=28, bold=True, color=BLACK)
    rect(slide, Inches(0.7), Inches(1.55), Inches(11.9), Inches(0.02), MGRAY)

def foot(slide, n, total=16):
    rect(slide, 0, H - Inches(0.4), W, Inches(0.4), LGRAY)
    txt(slide, "WiFi-Based Fall Detection System  ·  2026",
        Inches(0.7), H - Inches(0.38), Inches(9), Inches(0.32),
        size=9, color=DGRAY)
    txt(slide, f"{n} / {total}", Inches(12.0), H - Inches(0.38), Inches(1.1), Inches(0.32),
        size=9, color=DGRAY, align=PP_ALIGN.RIGHT)

def chip(slide, x, y, label, color=BLUE):
    w = Inches(len(label) * 0.105 + 0.35)
    rect(slide, x, y, w, Inches(0.32), color)
    txt(slide, label, x + Inches(0.12), y + Inches(0.04), w, Inches(0.28),
        size=9, bold=True, color=WHITE)
    return w

# ════════════════════════════════════════════════════════
# 1 — Cover
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, Inches(0.06), H, BLACK)          # left black bar
rect(s, Inches(0.06), H * 0.55, W, H * 0.45, LGRAY)  # bottom gray band

txt(s, "Fall Detection System", Inches(0.7), Inches(1.4), Inches(11), Inches(1.0),
    size=46, bold=True, color=BLACK)
txt(s, "WiFi CSI  ·  LSTM  ·  Elderly Care",
    Inches(0.7), Inches(2.45), Inches(8), Inches(0.5),
    size=18, color=DGRAY)
rect(s, Inches(0.7), Inches(3.1), Inches(1.2), Inches(0.04), BLUE)

txt(s, "ชื่อทีม  ·  สาขาวิศวกรรมคอมพิวเตอร์  ·  2026",
    Inches(0.7), Inches(5.0), Inches(10), Inches(0.4),
    size=13, color=DGRAY)
txt(s, "No Camera  ·  No Wearable  ·  Real-time Alert",
    Inches(0.7), Inches(5.5), Inches(10), Inches(0.4),
    size=13, color=BLUE, bold=True)
foot(s, 1)

# ════════════════════════════════════════════════════════
# 2 — Problem
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "PROBLEM", "ทำไมการตรวจจับการล้มจึงสำคัญ?")

txt(s, "684,000", Inches(0.7), Inches(1.7), Inches(5), Inches(1.6),
    size=80, bold=True, color=BLACK)
txt(s, "fall-related deaths / year  (WHO)",
    Inches(0.7), Inches(3.3), Inches(5), Inches(0.45),
    size=14, color=DGRAY)
rect(s, Inches(0.7), Inches(3.85), Inches(4.5), Inches(0.02), MGRAY)
txt(s, "ทุก 11 วินาที ผู้สูงอายุ 1 คนเข้า ER จากการล้ม",
    Inches(0.7), Inches(4.0), Inches(5.2), Inches(0.5),
    size=13, color=DGRAY, italic=True)

pains = [
    ("ตรวจจับล้มได้ช้า",     "ผู้สูงอายุอาจนอนอยู่คนเดียวนานหลายชั่วโมง"),
    ("ไม่มีใครช่วยทันที",    "Golden hour ของ Hip fracture คือ 6 ชั่วโมงแรก"),
    ("ผู้ดูแลเฝ้าไม่ได้ตลอด","ต้องการระบบอัตโนมัติที่เชื่อถือได้"),
]
for i, (t, d) in enumerate(pains):
    cy = Inches(1.8) + i * Inches(1.6)
    rect(s, Inches(6.5), cy, Inches(6.4), Inches(1.4), LGRAY, MGRAY)
    txt(s, t, Inches(6.75), cy + Inches(0.2), Inches(5.9), Inches(0.45),
        size=15, bold=True, color=BLACK)
    txt(s, d, Inches(6.75), cy + Inches(0.65), Inches(5.9), Inches(0.6),
        size=12, color=DGRAY)

foot(s, 2)

# ════════════════════════════════════════════════════════
# 3 — Existing Solutions
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "EXISTING SOLUTIONS", "ทางเลือกที่มีอยู่ — และข้อจำกัด")

cols = [
    ("กล้อง CCTV",       ["✗  ละเมิด Privacy", "✗  ต้องมีคนเฝ้า monitor", "✗  ไม่ครอบคลุมทุกมุม"]),
    ("Wearable / Pendant",["✗  ต้องสวมใส่ตลอดเวลา", "✗  ต้องชาร์จแบต", "✗  ผู้สูงอายุมักปฏิเสธ"]),
    ("ปุ่มกด SOS",        ["✗  ต้องกดเอง", "✗  กดไม่ได้ถ้าหมดสติ", "✗  ไม่ตรวจจับอัตโนมัติ"]),
    ("AI + Camera",       ["✗  ยังติด Privacy issue", "✗  ต้องการ GPU", "✗  แสงน้อยมีปัญหา"]),
]
for i, (name, cons) in enumerate(cols):
    cx = Inches(0.7) + i * Inches(3.15)
    rect(s, cx, Inches(1.75), Inches(2.9), Inches(4.8), LGRAY, MGRAY)
    txt(s, name, cx + Inches(0.2), Inches(1.9), Inches(2.5), Inches(0.55),
        size=14, bold=True, color=BLACK)
    rect(s, cx + Inches(0.2), Inches(2.45), Inches(1.5), Inches(0.02), MGRAY)
    for j, c in enumerate(cons):
        txt(s, c, cx + Inches(0.2), Inches(2.6) + j * Inches(0.65),
            Inches(2.5), Inches(0.6), size=12, color=RED)

rect(s, Inches(0.7), Inches(6.75), Inches(11.9), Inches(0.02), BLUE)
txt(s, "→  ยังไม่มีวิธีที่ passive · ไม่รุกล้ำ · ราคาถูก · ตรวจจับอัตโนมัติ",
    Inches(0.7), Inches(6.8), Inches(11.9), Inches(0.4),
    size=13, bold=True, color=BLUE)
foot(s, 3)

# ════════════════════════════════════════════════════════
# 4 — Our Solution
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, W, H, LGRAY)  # full gray bg for contrast slide
rect(s, 0, 0, Inches(6.5), H, BLACK)  # left black half

txt(s, "OUR SOLUTION", Inches(0.7), Inches(0.9), Inches(5.5), Inches(0.4),
    size=11, bold=True, color=BLUE)
txt(s, "ตรวจจับการล้ม\nผ่านสัญญาณ WiFi",
    Inches(0.7), Inches(1.35), Inches(5.5), Inches(2.0),
    size=36, bold=True, color=WHITE)
txt(s, "No Camera  ·  No Wearable\nReal-time Alert  ·  Auto Emergency Call",
    Inches(0.7), Inches(3.5), Inches(5.5), Inches(0.8),
    size=14, color=BLUE)

pillars = [
    ("Passive", "ใช้ WiFi ที่มีอยู่แล้ว\nไม่ต้องติดตั้งอุปกรณ์เพิ่ม"),
    ("AI-Powered", "LSTM  ·  97.9% Accuracy\nFall Recall 98.5%"),
    ("Auto Alert", "Push → Acknowledge\nหรือโทรฉุกเฉิน Twilio"),
]
for i, (title, desc) in enumerate(pillars):
    cy = Inches(1.75) + i * Inches(1.6)
    rect(s, Inches(7.0), cy, Inches(5.8), Inches(1.45), WHITE, MGRAY)
    rect(s, Inches(7.0), cy, Inches(0.07), Inches(1.45), BLUE)
    txt(s, title, Inches(7.25), cy + Inches(0.18), Inches(5.3), Inches(0.45),
        size=16, bold=True, color=BLACK)
    txt(s, desc, Inches(7.25), cy + Inches(0.65), Inches(5.3), Inches(0.65),
        size=12, color=DGRAY)

foot(s, 4)

# ════════════════════════════════════════════════════════
# 5 — WiFi CSI
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "TECHNOLOGY", "WiFi CSI — Channel State Information")

txt(s, "คลื่น WiFi เปลี่ยนรูปแบบเมื่อมีร่างกายเคลื่อนไหว\nระบบอ่านการเปลี่ยนแปลงนั้นเพื่อตรวจจับการล้ม",
    Inches(0.7), Inches(1.75), Inches(5.8), Inches(1.2),
    size=15, color=BLACK)

points = [
    "ทำงานในความมืดได้ — ไม่ต้องใช้แสง",
    "ทะลุผ่านกำแพงบางได้",
    "ไม่บันทึกภาพใดๆ — Privacy 100%",
    "ใช้ WiFi chip ESP32 ราคาถูก",
    "52 subcarrier  ·  100 packets/sec",
]
for i, p in enumerate(points):
    rect(s, Inches(0.7), Inches(3.1) + i * Inches(0.62), Inches(0.06), Inches(0.42), BLUE)
    txt(s, p, Inches(0.9), Inches(3.12) + i * Inches(0.62), Inches(5.3), Inches(0.42),
        size=13, color=BLACK)

# Signal diagram (right)
rect(s, Inches(7.0), Inches(1.75), Inches(5.8), Inches(2.3), LGRAY, MGRAY)
txt(s, "Normal Activity", Inches(7.2), Inches(1.85), Inches(5.4), Inches(0.35),
    size=11, bold=True, color=DGRAY)
txt(s, "Amplitude  ∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿∿",
    Inches(7.2), Inches(2.25), Inches(5.4), Inches(0.4), size=11, color=DGRAY)

rect(s, Inches(7.0), Inches(4.25), Inches(5.8), Inches(2.3), LGRAY, RED)
rect(s, Inches(7.0), Inches(4.25), Inches(5.8), Inches(0.04), RED)
txt(s, "FALL DETECTED", Inches(7.2), Inches(4.35), Inches(5.4), Inches(0.35),
    size=11, bold=True, color=RED)
txt(s, "Amplitude  ∿∿∿∿╔══╗∿∿╗\n               ╚══════╗∿∿∿",
    Inches(7.2), Inches(4.75), Inches(5.4), Inches(0.9), size=11, color=RED)

foot(s, 5)

# ════════════════════════════════════════════════════════
# 6 — Architecture Diagram
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ARCHITECTURE", "System Architecture — Input / Output / Data Flow")

# ── Diagram helpers ──────────────────────────────────────

def cbox(sl, x, y, w, h, hdr, body, hc=BLACK, bc=LGRAY):
    """Component box: colored header strip + body text."""
    rect(sl, x, y, w, h, bc, MGRAY)
    rect(sl, x, y, w, Inches(0.42), hc)
    txt(sl, hdr, x+Inches(0.13), y+Inches(0.07), w-Inches(0.16), Inches(0.32),
        size=11, bold=True, color=WHITE)
    txt(sl, body, x+Inches(0.13), y+Inches(0.5), w-Inches(0.16), h-Inches(0.58),
        size=10, color=DGRAY)

def h_arr(sl, x1, y, x2, label="", lc=BLUE):
    """Horizontal arrow x1→x2 at height y with optional label."""
    my = y - Inches(0.015)
    rect(sl, x1, my, x2-x1-Inches(0.12), Inches(0.03), lc)
    txt(sl, "▶", x2-Inches(0.18), y-Inches(0.14), Inches(0.2), Inches(0.28),
        size=9, bold=True, color=lc)
    if label:
        txt(sl, label, (x1+x2)/2-Inches(0.9), y-Inches(0.32), Inches(1.8), Inches(0.28),
            size=8, italic=True, color=lc, align=PP_ALIGN.CENTER)

def v_arr(sl, x, y1, y2, label="", lc=BLUE):
    """Vertical arrow y1→y2 at x with optional label."""
    rect(sl, x-Inches(0.015), y1, Inches(0.03), y2-y1-Inches(0.12), lc)
    txt(sl, "▼", x-Inches(0.1), y2-Inches(0.2), Inches(0.22), Inches(0.28),
        size=9, bold=True, color=lc)
    if label:
        txt(sl, label, x+Inches(0.08), (y1+y2)/2-Inches(0.12), Inches(1.4), Inches(0.28),
            size=8, italic=True, color=lc)

def sec_label(sl, x, y, label):
    txt(sl, label, x, y, Inches(3.0), Inches(0.28),
        size=8, bold=True, color=DGRAY)

# ── Column X positions ────────────────────────────────────
X_HW   = Inches(0.25)   # Hardware column
X_API  = Inches(4.35)   # Backend API column
X_ML   = Inches(8.0)    # ML Service column
X_MOB  = Inches(11.1)   # Mobile App column

BW = Inches(3.6)   # box width (hardware / API / ML)
MW = Inches(2.0)   # mobile box width

# ── Section labels ────────────────────────────────────────
sec_label(s, X_HW,  Inches(1.68), "① HARDWARE")
sec_label(s, X_API, Inches(1.68), "② BACKEND API  —  Render Cloud")
sec_label(s, X_ML,  Inches(1.68), "③ ML SERVICE  —  Render Cloud")
sec_label(s, X_MOB, Inches(1.68), "④ CLIENT")

# Thin section dividers
for xd in [X_API-Inches(0.18), X_ML-Inches(0.18), X_MOB-Inches(0.18)]:
    rect(s, xd, Inches(1.65), Inches(0.02), Inches(5.3), MGRAY)

# ── Hardware boxes ────────────────────────────────────────
# ESP32 #2 (STA)
cbox(s, X_HW, Inches(1.95), BW, Inches(1.7),
     "ESP32  #2  —  STA",
     "Station Mode\n"
     "ส่ง WiFi packet → ESP32 #1\n"
     "100 pkt/sec  ·  Channel 6",
     hc=RGBColor(0x33,0x33,0x33))

# WiFi CSI vertical arrow
v_arr(s, X_HW+BW/2, Inches(3.65), Inches(4.05), "WiFi CSI\nSignal", lc=DGRAY)

# ESP32 #1 (AP)
cbox(s, X_HW, Inches(4.05), BW, Inches(2.0),
     "ESP32  #1  —  AP",
     "Access Point Mode\n"
     "วัด CSI จาก ESP32 #2\n"
     "52 subcarriers × 100 pkt/s\n"
     "→ แยก 416 features/window",
     hc=BLACK)

# INPUT label
rect(s, X_HW, Inches(6.2), BW, Inches(0.55), RGBColor(0xEE,0xF2,0xF8), MGRAY)
txt(s, "INPUT   CSI raw · 52 subcarrier · 100 pkt/s",
    X_HW+Inches(0.12), Inches(6.28), BW-Inches(0.15), Inches(0.38),
    size=9, bold=True, color=BLUE)

# Arrow HW → API
h_arr(s, X_HW+BW, Inches(5.05), X_API,
      "HTTP POST\nCSI features", lc=BLUE)

# ── Backend API box ───────────────────────────────────────
cbox(s, X_API, Inches(1.95), BW, Inches(4.2),
     "Backend API  —  Node.js + Express",
     "• รับ CSI features จาก ESP32\n"
     "• เรียก ML Service → prediction\n"
     "• บันทึก event → Supabase DB\n"
     "• ส่ง Push → FCM\n"
     "• โทรฉุกเฉิน → Twilio\n"
     "• Acknowledge / Escalation logic",
     hc=BLUE, bc=RGBColor(0xF0,0xF5,0xFF))

# Service chips inside API box
for ci, (chip_label, cc) in enumerate([
    ("Supabase DB", RGBColor(0x3E,0xCF,0x8E)),
    ("Twilio",      RGBColor(0xF2,0x2F,0x46)),
    ("FCM Push",    RGBColor(0xFF,0x9A,0x00)),
]):
    chip_x = X_API + Inches(0.15) + ci * Inches(1.15)
    rect(s, chip_x, Inches(5.6), Inches(1.05), Inches(0.32), cc)
    txt(s, chip_label, chip_x+Inches(0.07), Inches(5.63), Inches(0.95), Inches(0.26),
        size=8, bold=True, color=WHITE)

# PROCESS label
rect(s, X_API, Inches(6.2), BW, Inches(0.55), RGBColor(0xEE,0xF2,0xFF), MGRAY)
txt(s, "PROCESS   predict() · confidence · risk_score",
    X_API+Inches(0.12), Inches(6.28), BW-Inches(0.15), Inches(0.38),
    size=9, bold=True, color=BLUE)

# ── Bidirectional arrow API ↔ ML ─────────────────────────
# Forward →
h_arr(s, X_API+BW, Inches(2.7), X_ML, "features [ ]", lc=BLUE)
# Return ←
rect(s, X_API+BW, Inches(3.1)-Inches(0.015), X_ML-(X_API+BW)-Inches(0.12), Inches(0.03),
     RGBColor(0x28,0x9E,0x8B))
txt(s, "◀", X_API+BW-Inches(0.05), Inches(3.1)-Inches(0.14), Inches(0.2), Inches(0.28),
    size=9, bold=True, color=RGBColor(0x28,0x9E,0x8B))
txt(s, "prediction · confidence",
    X_API+BW+Inches(0.1), Inches(3.1)+Inches(0.04), Inches(1.8), Inches(0.28),
    size=8, italic=True, color=RGBColor(0x28,0x9E,0x8B))

# ── ML Service box ────────────────────────────────────────
cbox(s, X_ML, Inches(1.95), BW, Inches(4.2),
     "ML Service  —  FastAPI  /  Python",
     "• รับ 416 features (10-window seq)\n"
     "• LSTM 128 → 64 → 32 units\n"
     "• Output: fall / no_fall\n"
     "• Confidence score\n"
     "• Risk score\n"
     "• Latency < 50 ms",
     hc=RGBColor(0x1A,0x53,0x96), bc=RGBColor(0xF0,0xF4,0xFF))

# Accuracy badge inside ML box
rect(s, X_ML+Inches(0.15), Inches(5.55), Inches(1.6), Inches(0.38), RGBColor(0x1A,0x53,0x96))
txt(s, "Accuracy  97.9%", X_ML+Inches(0.22), Inches(5.6), Inches(1.5), Inches(0.3),
    size=9, bold=True, color=WHITE)
rect(s, X_ML+Inches(1.85), Inches(5.55), Inches(1.6), Inches(0.38), RED)
txt(s, "Recall  98.5%", X_ML+Inches(1.92), Inches(5.6), Inches(1.5), Inches(0.3),
    size=9, bold=True, color=WHITE)

# ML PROCESS label
rect(s, X_ML, Inches(6.2), BW, Inches(0.55), RGBColor(0xEE,0xF2,0xFF), MGRAY)
txt(s, "MODEL   LSTM v3  ·  lstm_v3.h5  ·  scaler_v3.pkl",
    X_ML+Inches(0.12), Inches(6.28), BW-Inches(0.15), Inches(0.38),
    size=9, bold=True, color=RGBColor(0x1A,0x53,0x96))

# Arrow API → Mobile (FCM)
h_arr(s, X_ML+BW, Inches(3.5), X_MOB, "FCM Push\n< 1 sec", lc=RGBColor(0xFF,0x9A,0x00))

# ── Mobile App box ────────────────────────────────────────
cbox(s, X_MOB, Inches(1.95), MW, Inches(4.2),
     "Mobile App",
     "React Native\n+ Expo\n\n"
     "• รับ notification\n"
     "• Acknowledge\n"
     "• Event Log\n"
     "• ดู history",
     hc=RGBColor(0x28,0x9E,0x8B), bc=RGBColor(0xF0,0xFB,0xF7))

# OUTPUT label
rect(s, X_MOB, Inches(6.2), MW, Inches(0.55), RGBColor(0xEE,0xFB,0xF5), MGRAY)
txt(s, "OUTPUT   Alert + Call",
    X_MOB+Inches(0.12), Inches(6.28), MW-Inches(0.1), Inches(0.38),
    size=9, bold=True, color=RGBColor(0x28,0x9E,0x8B))

foot(s, 6)

# ════════════════════════════════════════════════════════
# 7 — Use Cases
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "USE CASES", "Scenario ที่ระบบต้องรับมือได้")

cases = [
    (BLUE, "Normal",
     "ผู้สูงอายุล้มในห้องที่ ESP32 ครอบคลุม",
     ["CSI pattern เปลี่ยนชัดเจน", "Confidence สูง > threshold", "Alert ส่งภายใน 1 วินาที"]),
    (RGBColor(0xF4,0xA2,0x61), "Edge Case",
     "กรณีที่เกิดได้ยาก แต่ต้องรับมือ",
     ["ล้มใกล้ขอบพื้นที่ → signal อ่อน confidence ต่ำ",
      "ของล้ม / คนเดินเร็ว → อาจ false positive",
      "ห้องน้ำ / กำแพงกั้น → ต้องเพิ่ม node"]),
    (RGBColor(0x28,0x9E,0x8B), "Business",
     "การแก้ปัญหาธุรกิจจริง",
     ["บ้านพักผู้สูงอายุ 30 ห้อง ติด ESP32 < ฿15,000",
      "ลด liability risk จากการช่วยเหลือล่าช้า",
      "ประหยัดค่าพนักงานเฝ้าตลอดคืนหลักแสน/ปี"]),
]
for i, (color, tag, title, bullets) in enumerate(cases):
    cx = Inches(0.7) + i * Inches(4.1)
    rect(s, cx, Inches(1.75), Inches(3.8), Inches(5.3), LGRAY, MGRAY)
    rect(s, cx, Inches(1.75), Inches(3.8), Inches(0.06), color)
    rect(s, cx, Inches(1.75), Inches(1.1), Inches(0.45), color)
    txt(s, tag, cx + Inches(0.08), Inches(1.8), Inches(1.0), Inches(0.35),
        size=10, bold=True, color=WHITE)
    txt(s, title, cx + Inches(0.2), Inches(2.3), Inches(3.4), Inches(0.55),
        size=14, bold=True, color=BLACK)
    rect(s, cx + Inches(0.2), Inches(2.9), Inches(2.0), Inches(0.02), MGRAY)
    for j, b in enumerate(bullets):
        rect(s, cx + Inches(0.2), Inches(3.05) + j * Inches(0.72), Inches(0.05), Inches(0.5), color)
        txt(s, b, cx + Inches(0.38), Inches(3.07) + j * Inches(0.72),
            Inches(3.2), Inches(0.65), size=12, color=DGRAY)

foot(s, 7)

# ════════════════════════════════════════════════════════
# 8 — Data Pipeline
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "DATA PIPELINE", "จากสัญญาณ ESP32 → โมเดล AI")

steps = [
    ("01", "Collect",     "ESP32 → USB Serial\n400 pkts / file\nCSV raw data"),
    ("02", "Preprocess",  "Window 200 · Stride 50\n8 stats × 52 sub\n= 416 features"),
    ("03", "Train",       "LSTM 128→64→32\n1,242 windows\ntrain_v3.ipynb"),
    ("04", "Inference",   "FastAPI · Render\nReal-time < 50ms\nconfidence + risk"),
]
for i, (num, title, desc) in enumerate(steps):
    cx = Inches(0.7) + i * Inches(3.15)
    rect(s, cx, Inches(1.75), Inches(2.85), Inches(5.2), LGRAY, MGRAY)
    txt(s, num, cx + Inches(0.2), Inches(1.9), Inches(1.0), Inches(0.8),
        size=36, bold=True, color=BLUE)
    rect(s, cx + Inches(0.2), Inches(2.7), Inches(2.1), Inches(0.02), BLUE)
    txt(s, title, cx + Inches(0.2), Inches(2.85), Inches(2.5), Inches(0.45),
        size=16, bold=True, color=BLACK)
    txt(s, desc, cx + Inches(0.2), Inches(3.4), Inches(2.5), Inches(2.0),
        size=12, color=DGRAY)
    if i < 3:
        txt(s, "→", cx + Inches(2.85), Inches(4.0), Inches(0.4), Inches(0.5),
            size=22, color=MGRAY, align=PP_ALIGN.CENTER)

foot(s, 8)

# ════════════════════════════════════════════════════════
# 8 — AI Model
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "AI MODEL", "LSTM Architecture & Performance")

# Architecture (left)
arch = [
    ("Input",   "(10, 416)  — Sequence × Features"),
    ("LSTM 1",  "128 units  ·  return_sequences=True"),
    ("LSTM 2",  " 64 units  ·  return_sequences=True"),
    ("LSTM 3",  " 32 units"),
    ("Output",  "2 classes  ·  Softmax  →  fall / no_fall"),
]
for i, (layer, detail) in enumerate(arch):
    cy = Inches(1.75) + i * Inches(0.85)
    fill = BLACK if i in [0, 4] else LGRAY
    lcolor = WHITE if i in [0, 4] else BLACK
    dcolor = BLUE if i in [0, 4] else DGRAY
    rect(s, Inches(0.7), cy, Inches(5.8), Inches(0.75), fill, MGRAY)
    txt(s, layer, Inches(0.9), cy + Inches(0.17), Inches(1.2), Inches(0.42),
        size=12, bold=True, color=lcolor)
    txt(s, detail, Inches(2.2), cy + Inches(0.17), Inches(4.1), Inches(0.42),
        size=12, color=dcolor)

txt(s, "416 features  ·  10 sequence steps  ·  Model v3",
    Inches(0.7), Inches(6.2), Inches(5.8), Inches(0.35),
    size=10, color=DGRAY, italic=True)

# Metrics (right)
metrics = [
    ("97.9%", "Overall Accuracy", BLUE),
    ("98.5%", "Fall Recall",      RED),
    ("97.4%", "Precision",        BLACK),
    ("1,242", "Training Windows", DGRAY),
]
for i, (val, label, color) in enumerate(metrics):
    cx = Inches(7.2) + (i % 2) * Inches(3.0)
    cy = Inches(1.75) + (i // 2) * Inches(2.55)
    rect(s, cx, cy, Inches(2.7), Inches(2.3), LGRAY, MGRAY)
    txt(s, val, cx + Inches(0.15), cy + Inches(0.3), Inches(2.4), Inches(1.1),
        size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, label, cx + Inches(0.15), cy + Inches(1.5), Inches(2.4), Inches(0.5),
        size=12, color=DGRAY, align=PP_ALIGN.CENTER)

foot(s, 9)

# ════════════════════════════════════════════════════════
# 9 — Dataset
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "DATASET", "ข้อมูลที่เก็บและใช้ Train โมเดล")

rows = [
    ("fall_A",   "100", "300", "การล้มไปด้านหน้า"),
    ("fall_B",   "100", "300", "การล้มไปด้านข้าง"),
    ("fall_C",   "100", "300", "การล้มไปด้านหลัง"),
    ("non_fall", "120", "360", "เดิน · ยืน · นั่ง · นอน"),
]
headers = ["Class", "ไฟล์", "Windows", "สัดส่วน", "คำอธิบาย"]
cxs = [Inches(0.7), Inches(3.0), Inches(5.0), Inches(7.0), Inches(8.8)]
cws = [Inches(2.1), Inches(1.8), Inches(1.8), Inches(1.6), Inches(3.8)]

rect(s, Inches(0.7), Inches(1.75), Inches(12.0), Inches(0.55), BLACK)
for hdr, cx, cw in zip(headers, cxs, cws):
    txt(s, hdr, cx + Inches(0.1), Inches(1.82), cw, Inches(0.42),
        size=12, bold=True, color=WHITE)

total_w = sum(int(r[2]) for r in rows)
for i, (cls, files, wins, desc) in enumerate(rows):
    cy = Inches(2.3) + i * Inches(0.95)
    fill = WHITE if i % 2 == 0 else LGRAY
    rect(s, Inches(0.7), cy, Inches(12.0), Inches(0.9), fill, MGRAY)
    pct = f"{int(wins)/total_w*100:.0f}%"
    vals = [cls, files, wins, pct, desc]
    for j, (val, cx, cw) in enumerate(zip(vals, cxs, cws)):
        c = RED if j == 0 and "fall" in cls and cls != "non_fall" else \
            RGBColor(0x28,0x9E,0x8B) if j == 0 and cls == "non_fall" else BLACK
        txt(s, val, cx + Inches(0.1), cy + Inches(0.22), cw, Inches(0.45),
            size=13, bold=(j == 0), color=c)

cy = Inches(2.3) + 4 * Inches(0.95)
rect(s, Inches(0.7), cy, Inches(12.0), Inches(0.9), BLACK)
for val, cx, cw in zip(["รวม", "420", "1,260", "100%", "—"], cxs, cws):
    txt(s, val, cx + Inches(0.1), cy + Inches(0.22), cw, Inches(0.45),
        size=13, bold=True, color=WHITE)

foot(s, 10)

# ════════════════════════════════════════════════════════
# 10 — Alert Flow
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "ALERT FLOW", "ระบบแจ้งเตือนอัตโนมัติ")

nodes = [
    (RED,                        "FALL\nDETECTED",     "ML confidence\n> threshold"),
    (BLUE,                       "PUSH\nNOTIFICATION", "FCM → แอปผู้ดูแล\n< 1 วินาที"),
    (RGBColor(0xF4,0xA2,0x61),  "WAIT\n30 sec",       "รอ Acknowledge\nจากผู้ดูแล"),
    (RGBColor(0x28,0x9E,0x8B),  "ACKNOWLEDGED",        "ผู้ดูแลกด OK\nระบบหยุดแจ้งเตือน"),
]
for i, (color, title, desc) in enumerate(nodes):
    cx = Inches(0.7) + i * Inches(2.9)
    rect(s, cx, Inches(1.85), Inches(2.6), Inches(3.2), color)
    txt(s, title, cx + Inches(0.1), Inches(2.1), Inches(2.4), Inches(1.0),
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, desc, cx + Inches(0.1), Inches(3.2), Inches(2.4), Inches(0.8),
        size=12, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(s, "→", cx + Inches(2.6), Inches(3.0), Inches(0.4), Inches(0.5),
            size=20, color=DGRAY, align=PP_ALIGN.CENTER)

# Escalation
rect(s, Inches(0.7), Inches(5.3), Inches(0.06), Inches(1.5), RED)
rect(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(0.02), MGRAY)
txt(s, "ถ้าไม่มีใครตอบภายใน 30 วินาที  →",
    Inches(0.9), Inches(5.35), Inches(5.0), Inches(0.4),
    size=13, color=BLACK)
rect(s, Inches(6.1), Inches(5.25), Inches(6.5), Inches(1.3), RED)
txt(s, "Twilio Emergency Call  —  โทรฉุกเฉินอัตโนมัติ",
    Inches(6.3), Inches(5.55), Inches(6.1), Inches(0.4),
    size=14, bold=True, color=WHITE)

foot(s, 11)

# ════════════════════════════════════════════════════════
# 11 — DEMO (video placeholder)
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, W, H, BLACK)

txt(s, "LIVE DEMO", Inches(0.7), Inches(0.5), Inches(6), Inches(0.4),
    size=11, bold=True, color=BLUE)
txt(s, "การสาธิตระบบ", Inches(0.7), Inches(0.9), Inches(6), Inches(0.7),
    size=32, bold=True, color=WHITE)
rect(s, Inches(0.7), Inches(1.6), Inches(1.0), Inches(0.04), BLUE)

# Video placeholder box
rect(s, Inches(0.7), Inches(1.8), Inches(8.5), Inches(5.25), RGBColor(0x1A,0x1A,0x1A))
# dashed border simulation with thin blue rect outline
for side_x, side_y, side_w, side_h in [
    (Inches(0.7),  Inches(1.8),  Inches(8.5), Inches(0.04)),   # top
    (Inches(0.7),  Inches(7.0),  Inches(8.5), Inches(0.04)),   # bottom
    (Inches(0.7),  Inches(1.8),  Inches(0.04), Inches(5.25)),  # left
    (Inches(9.16), Inches(1.8),  Inches(0.04), Inches(5.25)),  # right
]:
    rect(s, side_x, side_y, side_w, side_h, BLUE)

txt(s, "▶  แทรก Video ที่นี่\nInsert > Video > Video on My PC…",
    Inches(2.5), Inches(4.0), Inches(4.5), Inches(1.0),
    size=16, color=DGRAY, align=PP_ALIGN.CENTER)

# Side notes
notes = [
    ("สิ่งที่จะเห็นในคลิป", True),
    ("ESP32 ตั้งอยู่ในห้อง", False),
    ("ผู้ทดสอบล้มในพื้นที่", False),
    ("แอปรับ notification < 1s", False),
    ("ผู้ดูแล Acknowledge บนมือถือ", False),
]
for i, (note, bold) in enumerate(notes):
    color = WHITE if bold else DGRAY
    size  = 13 if bold else 12
    if not bold:
        rect(s, Inches(9.5), Inches(2.5) + i * Inches(0.62), Inches(0.05), Inches(0.45), BLUE)
    txt(s, note, Inches(9.65) if not bold else Inches(9.4),
        Inches(2.48) + i * Inches(0.62), Inches(3.5), Inches(0.5),
        size=size, bold=bold, color=color)

foot(s, 12)

# ════════════════════════════════════════════════════════
# 12 — Demo & Results
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "DEMO & RESULTS", "การทดสอบระบบ End-to-End")

steps_demo = [
    "ESP32 วางในห้อง — สร้าง WiFi และส่ง CSI",
    "ผู้ทดสอบล้ม — CSI pattern เปลี่ยน",
    "ML Service บน Render Cloud ทำนาย < 50ms",
    "Express API บันทึก Supabase + ส่ง FCM",
    "แอปรับ notification < 1 วินาที",
    "กด Acknowledge หรือรอรับสาย Twilio",
]
for i, step in enumerate(steps_demo):
    rect(s, Inches(0.7), Inches(1.75) + i * Inches(0.72), Inches(0.06), Inches(0.52), BLUE)
    txt(s, step, Inches(0.95), Inches(1.77) + i * Inches(0.72), Inches(5.4), Inches(0.52),
        size=13, color=BLACK)

kpis = [
    ("< 1s",  "Detection\nLatency",   BLUE),
    ("97.9%", "Accuracy",             BLACK),
    ("98.5%", "Fall Recall",          RED),
    ("~฿350", "Hardware\nCost/Unit",  RGBColor(0x28,0x9E,0x8B)),
]
for i, (val, label, color) in enumerate(kpis):
    cx = Inches(6.8) + (i % 2) * Inches(3.3)
    cy = Inches(1.75) + (i // 2) * Inches(2.65)
    rect(s, cx, cy, Inches(3.0), Inches(2.4), LGRAY, MGRAY)
    rect(s, cx, cy, Inches(3.0), Inches(0.06), color)
    txt(s, val, cx + Inches(0.15), cy + Inches(0.3), Inches(2.7), Inches(1.1),
        size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(s, label, cx + Inches(0.15), cy + Inches(1.5), Inches(2.7), Inches(0.7),
        size=12, color=DGRAY, align=PP_ALIGN.CENTER)

foot(s, 13)

# ════════════════════════════════════════════════════════
# 12 — Impact
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "IMPACT", "กลุ่มเป้าหมายและการเปรียบเทียบ")

targets = [
    ("บ้านพักผู้สูงอายุ",     "วาง ESP32 node หลายตัว\nครอบคลุมทุกห้อง ราคาถูก"),
    ("ครอบครัว",               "ติดตั้งเองได้\nลูกหลานรับแจ้งเตือนทางแอป"),
    ("โรงพยาบาล / Clinic",    "ใช้ WiFi ที่มีอยู่\nไม่ต้องลงทุนโครงสร้างใหม่"),
]
for i, (name, desc) in enumerate(targets):
    cx = Inches(0.7) + i * Inches(4.1)
    rect(s, cx, Inches(1.75), Inches(3.8), Inches(2.4), LGRAY, MGRAY)
    rect(s, cx, Inches(1.75), Inches(0.06), Inches(2.4), BLUE)
    txt(s, name, cx + Inches(0.2), Inches(1.9), Inches(3.4), Inches(0.5),
        size=15, bold=True, color=BLACK)
    txt(s, desc, cx + Inches(0.2), Inches(2.45), Inches(3.4), Inches(1.4),
        size=12, color=DGRAY)

# Cost table
rect(s, Inches(0.7), Inches(4.4), Inches(12.0), Inches(0.5), BLACK)
txt(s, "เปรียบเทียบต้นทุน", Inches(0.9), Inches(4.52), Inches(5), Inches(0.35),
    size=12, bold=True, color=WHITE)

costs = [
    ("ระบบนี้  (ESP32 ×2)", "~฿350–400",     RGBColor(0x28,0x9E,0x8B)),
    ("IP Camera",           "~฿1,500–5,000",  RED),
    ("Wearable Premium",    "~฿3,000–15,000", RED),
]
for i, (name, price, color) in enumerate(costs):
    cy = Inches(5.05) + i * Inches(0.55)
    fill = WHITE if i % 2 == 0 else LGRAY
    rect(s, Inches(0.7), cy, Inches(12.0), Inches(0.52), fill, MGRAY)
    txt(s, name, Inches(0.9), cy + Inches(0.1), Inches(5.0), Inches(0.35), size=13, color=BLACK)
    txt(s, price, Inches(6.0), cy + Inches(0.1), Inches(3.0), Inches(0.35),
        size=13, bold=True, color=color)

foot(s, 14)

# ════════════════════════════════════════════════════════
# 13 — Future Work
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, "FUTURE WORK", "แนวทางพัฒนาต่อ")

futures = [
    ("BLE Provisioning",       "ตั้งค่า WiFi บ้านผ่านแอป\nไม่ต้อง hardcode SSID"),
    ("One-Class Model",        "ฝึกโมเดลด้วย normal เท่านั้น\nตรวจจับความผิดปกติทุกรูปแบบ"),
    ("Multi-Room Coverage",    "ESP32 หลาย node\nระบุตำแหน่งการล้มได้"),
    ("Video Validation",       "บันทึก video คู่กับ CSI\nพิสูจน์ความแม่นยำ"),
    ("Larger Dataset",         "เพิ่ม distance variation\nผู้ทดสอบหลากหลายกลุ่ม"),
    ("Full Cloud Deployment",  "Deploy ครบ stack\nรองรับหลายบ้านพร้อมกัน"),
]
for i, (title, desc) in enumerate(futures):
    cx = Inches(0.7) + (i % 3) * Inches(4.1)
    cy = Inches(1.75) + (i // 3) * Inches(2.55)
    rect(s, cx, cy, Inches(3.8), Inches(2.3), LGRAY, MGRAY)
    rect(s, cx, cy, Inches(3.8), Inches(0.06), BLUE)
    txt(s, title, cx + Inches(0.2), Inches(0.2) + cy, Inches(3.4), Inches(0.5),
        size=14, bold=True, color=BLACK)
    txt(s, desc, cx + Inches(0.2), cy + Inches(0.75), Inches(3.4), Inches(1.3),
        size=12, color=DGRAY)

foot(s, 15)

# ════════════════════════════════════════════════════════
# 14 — Conclusion
# ════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
rect(s, 0, 0, W, H, BLACK)

txt(s, "CONCLUSION", Inches(0.7), Inches(0.7), Inches(10), Inches(0.4),
    size=11, bold=True, color=BLUE)
txt(s, "สรุป", Inches(0.7), Inches(1.1), Inches(10), Inches(0.7),
    size=38, bold=True, color=WHITE)
rect(s, Inches(0.7), Inches(1.8), Inches(1.2), Inches(0.04), BLUE)

points = [
    ("No Camera · No Wearable",  "ใช้ WiFi CSI ที่มีอยู่ — ติดตั้งง่าย ไม่รุกล้ำ Privacy"),
    ("AI Accuracy 97.9%",        "LSTM v3 · Fall Recall 98.5% · 1,242 training windows"),
    ("Alert < 1 วินาที",         "Push → Acknowledge → Twilio Emergency Call อัตโนมัติ"),
    ("Hardware ~฿350–400/จุด",   "ถูกกว่า IP Camera และ Wearable premium 10–40 เท่า"),
]
for i, (title, desc) in enumerate(points):
    cy = Inches(2.2) + i * Inches(1.15)
    rect(s, Inches(0.7), cy, Inches(12.0), Inches(1.0), RGBColor(0x1A,0x1A,0x1A))
    rect(s, Inches(0.7), cy, Inches(0.06), Inches(1.0),
         BLUE if i % 2 == 0 else RED)
    txt(s, title, Inches(0.95), cy + Inches(0.1), Inches(11.5), Inches(0.42),
        size=16, bold=True, color=WHITE)
    txt(s, desc, Inches(0.95), cy + Inches(0.55), Inches(11.5), Inches(0.38),
        size=12, color=DGRAY)

rect(s, 0, H - Inches(0.55), W, Inches(0.55), BLUE)
txt(s, "Thank you  ·  ขอบคุณครับ/ค่ะ",
    Inches(0.7), H - Inches(0.5), Inches(12.5), Inches(0.4),
    size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

foot(s, 16)

# ── Save ──────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "fall_detection_slides.pptx")
prs.save(out)
print(f"Saved → {out}")
