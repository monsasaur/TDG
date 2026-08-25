from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


_OUT_DIR = Path(__file__).resolve().parent.parent
OUT = str(_OUT_DIR / "tdg_lit_review_progress_30min.pptx")
NOTES_OUT = str(_OUT_DIR / "tdg_lit_review_progress_30min_notes.md")

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 25

BLACK = RGBColor(13, 17, 23)
INK = RGBColor(27, 35, 46)
MUTED = RGBColor(91, 103, 117)
SOFT = RGBColor(246, 248, 251)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(222, 229, 237)
RED = RGBColor(230, 57, 70)
RED_DARK = RGBColor(174, 35, 52)
BLUE = RGBColor(46, 102, 214)
TEAL = RGBColor(0, 142, 130)
GREEN = RGBColor(36, 150, 86)
AMBER = RGBColor(238, 153, 48)
PURPLE = RGBColor(111, 89, 204)

FONT = "Thonburi"
LATIN = "Aptos"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

notes = []


def rgb(value: str) -> RGBColor:
    value = value.replace("#", "")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16))


def rect(slide, x, y, w, h, fill=WHITE, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def rounded(slide, x, y, w, h, fill=WHITE, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def circle(slide, x, y, d, fill=WHITE, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.9)
    else:
        shape.line.fill.background()
    return shape


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    font=FONT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def center_text(slide, value, x, y, w, h, size=18, color=INK, bold=False, font=FONT):
    box = text(slide, value, x, y, w, h, size, color, bold, PP_ALIGN.CENTER, font)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def line(slide, x1, y1, x2, y2, color=LINE, width=2):
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def bg(slide):
    rect(slide, 0, 0, W, H, SOFT)


def kicker(slide, value, n):
    text(slide, value.upper(), Inches(0.62), Inches(0.32), Inches(7), Inches(0.25), 9, RED, True, font=LATIN)
    text(slide, f"{n:02d}/{TOTAL}", Inches(11.7), Inches(0.32), Inches(1.0), Inches(0.25), 9, MUTED, True, PP_ALIGN.RIGHT, LATIN)


def footer(slide, timing):
    rect(slide, Inches(0.62), Inches(7.03), Inches(12.08), Inches(0.01), LINE)
    text(slide, timing, Inches(10.45), Inches(7.11), Inches(2.25), Inches(0.25), 8.5, MUTED, True, PP_ALIGN.RIGHT, LATIN)


def headline(slide, value, sub=None, size=34):
    text(slide, value, Inches(0.62), Inches(0.62), Inches(12.1), Inches(0.9), size, INK, True)
    if sub:
        text(slide, sub, Inches(0.66), Inches(1.5), Inches(11.6), Inches(0.42), 15, MUTED)


def source(slide, value):
    text(slide, f"ที่มา: {value}", Inches(0.66), Inches(6.72), Inches(12.0), Inches(0.22), 7.2, MUTED)


def add_note(n, title, timing, body):
    notes.append((n, title, timing, body))


def pill(slide, label, x, y, color, w=None):
    if w is None:
        w = Inches(max(1.1, len(label) * 0.11 + 0.35))
    rounded(slide, x, y, w, Inches(0.34), color, None)
    center_text(slide, label, x + Inches(0.06), y + Inches(0.05), w - Inches(0.12), Inches(0.2), 9, WHITE, True)
    return w


def card(slide, x, y, w, h, title, body, color=BLUE, title_size=14, body_size=11.2):
    rounded(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x, y, Inches(0.08), h, color)
    text(slide, title, x + Inches(0.22), y + Inches(0.17), w - Inches(0.36), Inches(0.32), title_size, INK, True)
    text(slide, body, x + Inches(0.22), y + Inches(0.56), w - Inches(0.42), h - Inches(0.7), body_size, MUTED)


def metric(slide, x, y, value, label, color=BLUE, w=2.25):
    rounded(slide, x, y, Inches(w), Inches(1.05), WHITE, LINE)
    center_text(slide, value, x, y + Inches(0.15), Inches(w), Inches(0.35), 23, color, True, LATIN)
    center_text(slide, label, x + Inches(0.1), y + Inches(0.57), Inches(w - 0.2), Inches(0.33), 9.5, MUTED)


def bullets(slide, items, x, y, w, size=14, gap=0.55, color=INK, accent=RED):
    for i, item in enumerate(items):
        yy = y + Inches(i * gap)
        rect(slide, x, yy + Inches(0.12), Inches(0.07), Inches(0.18), accent)
        text(slide, item, x + Inches(0.2), yy, w - Inches(0.2), Inches(0.42), size, color)


def flow_step(slide, x, y, label, body, color, w=1.7):
    circle(slide, x + Inches(0.35), y, Inches(1.0), color, None)
    center_text(slide, label, x + Inches(0.35), y + Inches(0.33), Inches(1.0), Inches(0.2), 11, WHITE, True, LATIN)
    center_text(slide, body, x, y + Inches(1.18), Inches(w), Inches(0.48), 11.5, INK, True)


def table_cell(slide, x, y, w, h, value, fill=WHITE, color=INK, size=10.5, bold=False, align=PP_ALIGN.CENTER):
    rect(slide, x, y, w, h, fill, LINE)
    center_text(slide, value, x + Inches(0.04), y + Inches(0.04), w - Inches(0.08), h - Inches(0.08), size, color, bold)


# 1 Cover
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLACK)
rect(s, Inches(7.55), 0, Inches(5.8), H, rgb("#F2F5F9"))
kicker(s, "ภาพรวม", 1)
text(s, "Middle\nระบบตรวจจับการล้มด้วย WiFi CSI", Inches(0.72), Inches(0.95), Inches(6.5), Inches(1.45), 35, WHITE, True)
text(s, "ตั้งแต่ปัญหา งานวิจัยที่เกี่ยวข้อง ไปจนถึง prototype ที่ทำแล้ว", Inches(0.76), Inches(2.72), Inches(6.1), Inches(0.55), 16, rgb("#DCE5EF"))
rect(s, Inches(0.76), Inches(3.55), Inches(1.45), Inches(0.08), RED)
text(s, "เวลา 30 นาที รวมถามตอบ\nพรี 24 นาที + Q&A 6 นาที", Inches(0.76), Inches(4.0), Inches(5.7), Inches(0.9), 22, WHITE, True)
rounded(s, Inches(8.35), Inches(1.05), Inches(3.55), Inches(5.35), WHITE, LINE)
circle(s, Inches(9.45), Inches(1.55), Inches(1.25), rgb("#F7C8B6"), None)
line(s, Inches(10.08), Inches(2.8), Inches(10.08), Inches(4.0), INK, 4)
line(s, Inches(10.08), Inches(3.25), Inches(9.45), Inches(3.75), INK, 3)
line(s, Inches(10.08), Inches(3.25), Inches(10.72), Inches(3.75), INK, 3)
rounded(s, Inches(8.85), Inches(4.85), Inches(2.55), Inches(0.74), rgb("#FFE7EB"), None)
center_text(s, "FALL ALERT", Inches(9.05), Inches(5.04), Inches(2.15), Inches(0.27), 17, RED_DARK, True, LATIN)
text(s, "Middle / TDG Project", Inches(8.5), Inches(6.02), Inches(3.25), Inches(0.28), 12, MUTED, True, PP_ALIGN.CENTER, LATIN)
footer(s, "0:00 - 0:40")
add_note(1, "เปิดเรื่อง", "0:00 - 0:40", "เปิดด้วยโจทย์หลัก: ผู้สูงอายุล้มแล้วไม่มีใครรู้ทันที ระบบนี้พยายามตรวจจับแบบ passive โดยไม่ใช้กล้องและไม่ต้องใส่อุปกรณ์")


# 2 Timebox
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "ลำดับการนำเสนอ", 2)
headline(s, "วันนี้จะตอบ 4 คำถามหลัก", "ปัญหาคืออะไร ทำไมวิธีเดิมยังไม่พอ ทำไมเลือก WiFi CSI และตอนนี้ทำถึงไหนแล้ว")
segments = [
    ("1", "Pain point", "ทำไม fall detection สำคัญ", RED, "1:10-3:30"),
    ("2", "Lit review", "ข้อจำกัดของ camera / wearable / SOS และงาน WiFi CSI", BLUE, "3:30-12:20"),
    ("3", "Solution", "สถาปัตยกรรมระบบ Middle", TEAL, "12:20-18:00"),
    ("4", "Progress", "ทำอะไรเสร็จแล้ว, ยังเสี่ยงตรงไหน", GREEN, "18:00-24:00"),
    ("5", "Q&A", "ถามตอบและ references สำรอง", AMBER, "24:00-30:00"),
]
for i, (num, title, body, color, timing) in enumerate(segments):
    x = Inches(0.78 + i * 2.48)
    rounded(s, x, Inches(2.23), Inches(2.03), Inches(2.75), WHITE, LINE)
    circle(s, x + Inches(0.57), Inches(2.56), Inches(0.88), color, None)
    center_text(s, num, x + Inches(0.57), Inches(2.82), Inches(0.88), Inches(0.18), 12, WHITE, True, LATIN)
    center_text(s, title, x + Inches(0.14), Inches(3.65), Inches(1.75), Inches(0.3), 15, INK, True)
    center_text(s, body, x + Inches(0.16), Inches(4.05), Inches(1.72), Inches(0.52), 10.5, MUTED)
    center_text(s, timing, x + Inches(0.2), Inches(4.67), Inches(1.62), Inches(0.22), 8.5, color, True, LATIN)
    if i < 4:
        line(s, x + Inches(2.03), Inches(3.6), x + Inches(2.35), Inches(3.6), rgb("#B8C2CE"), 2)
rounded(s, Inches(1.18), Inches(5.7), Inches(10.95), Inches(0.72), BLACK, None)
center_text(s, "แกนกลางของคำตอบ: งานวิจัยบอกว่า WiFi CSI มี promise แต่ challenge คือ real-world generalization", Inches(1.45), Inches(5.88), Inches(10.4), Inches(0.26), 15.5, WHITE, True)
footer(s, "0:40 - 1:10")
add_note(2, "โครงเรื่องและเวลา", "0:40 - 1:10", "บอกกรรมการว่าจะไม่ลงลึกทุก implementation line แต่จะเชื่อม lit review กับการตัดสินใจทางระบบและสถานะที่ทำไปแล้ว")


# 3 Pain point stats
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "ปัญหา", 3)
headline(s, "ปัญหาหลักคือการรู้ช้าเมื่อผู้สูงอายุล้ม", "การล้มมีผลต่อชีวิต การรักษา และความกังวลของผู้ดูแล")
metric(s, Inches(0.86), Inches(2.0), "684k", "เสียชีวิตจากการล้มต่อปีทั่วโลก", RED, 2.45)
metric(s, Inches(3.65), Inches(2.0), "37.3M", "ล้มรุนแรงพอให้ต้องพบแพทย์ต่อปี", BLUE, 2.45)
metric(s, Inches(6.45), Inches(2.0), "1 in 4", "ผู้สูงอายุ 65+ ในสหรัฐฯ ล้มทุกปี", TEAL, 2.45)
metric(s, Inches(9.25), Inches(2.0), "3M", "ED visits ต่อปีจาก falls ในผู้สูงอายุ", AMBER, 2.45)
rounded(s, Inches(0.88), Inches(4.0), Inches(11.7), Inches(1.55), WHITE, LINE)
text(s, "Pain point ของ caregiver", Inches(1.18), Inches(4.28), Inches(3.3), Inches(0.3), 16, INK, True)
bullets(s, [
    "ถ้าไม่มีคนเห็น เหตุการณ์หนึ่งอาจกลายเป็น long lie และ delayed treatment",
    "ระบบที่ต้องให้ผู้สูงอายุกดเองมี failure mode สำคัญ: หมดสติ เจ็บ หรือเอื้อมไม่ถึง",
    "ครอบครัวต้องการ alert ที่เร็ว แต่ต้องลด false alarm จนไม่รบกวนชีวิตประจำวัน",
], Inches(1.18), Inches(4.72), Inches(10.8), 12.5, 0.34, MUTED, RED)
source(s, "WHO Falls fact sheet, 2021; CDC Facts About Falls, 2026")
footer(s, "1:10 - 2:20")
add_note(3, "ขนาดของปัญหา", "1:10 - 2:20", "ใช้ตัวเลขใหญ่เพื่อยืนยันว่า problem มีน้ำหนัก จากนั้นโยงกลับมาที่ pain point จริงของครอบครัว: รู้ช้าและช่วยช้า")


# 4 Requirements
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "โจทย์ออกแบบ", 4)
headline(s, "ระบบที่เหมาะกับบ้านต้องไม่เพิ่มภาระให้ผู้สูงอายุ", "จึงต้อง passive, privacy-first, real-time และมีแผนสำรองเมื่อไม่มีคนตอบรับ")
reqs = [
    ("Passive", "ผู้สูงอายุไม่ต้องกดปุ่ม ไม่ต้องจำใส่อุปกรณ์", RED),
    ("Privacy-first", "ไม่เก็บภาพหรือเสียงละเอียดในพื้นที่ส่วนตัว", BLUE),
    ("Real-time", "แจ้งเตือนภายในหลักวินาที ไม่ใช่หลังตรวจย้อนหลัง", TEAL),
    ("Escalation", "ถ้า caregiver ไม่ตอบรับ ต้องมี fallback ไปยังเบอร์ฉุกเฉิน", AMBER),
    ("Low-cost", "ใช้ hardware ราคาถูก ติดตั้งง่ายกว่าระบบ sensor หนาแน่น", GREEN),
]
for i, (title, body, color) in enumerate(reqs):
    x = Inches(0.78 + (i % 3) * 4.17)
    y = Inches(2.05 + (i // 3) * 1.72)
    w = Inches(3.55 if i < 3 else 5.55)
    card(s, x, y, w, Inches(1.25), title, body, color, 14, 11.3)
rounded(s, Inches(1.2), Inches(5.8), Inches(10.9), Inches(0.65), BLACK, None)
center_text(s, "โจทย์วิจัยของเรา: มีวิธี sensing ที่ passive และ privacy-first แต่ยังส่ง alert ได้จริงหรือไม่", Inches(1.45), Inches(5.98), Inches(10.4), Inches(0.24), 15, WHITE, True)
footer(s, "2:20 - 3:30")
add_note(4, "Requirement", "2:20 - 3:30", "พูดว่าข้อกำหนดเหล่านี้ทำให้ camera, wearable และ SOS button ไม่ใช่คำตอบเดียว จึงต้องดู literature ของ device-free sensing")


# 5 Existing solution taxonomy
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "วิธีเดิม", 5)
headline(s, "วิธีที่มีอยู่ช่วยได้ แต่ยังมีข้อแลกเปลี่ยน", "แต่ละกลุ่มแก้ปัญหาบางส่วน และสร้างข้อจำกัดใหม่เมื่อนำไปใช้จริง")
headers = ["กลุ่ม", "ตัวอย่าง", "จุดแข็ง", "ข้อจำกัดที่เจอบ่อย"]
widths = [1.65, 2.65, 3.0, 4.7]
x0 = Inches(0.72)
y0 = Inches(2.0)
row_h = Inches(0.65)
for i, h in enumerate(headers):
    table_cell(s, x0 + Inches(sum(widths[:i])), y0, Inches(widths[i]), row_h, h, BLACK, WHITE, 10.5, True)
rows = [
    ("Wearable", "watch, pendant, IMU", "แม่นใน controlled movement", "ต้องใส่และชาร์จ, placement sensitivity"),
    ("Vision", "RGB, depth camera", "เห็น posture ชัด", "privacy, occlusion, แสงและมุมกล้อง"),
    ("Ambient", "floor, acoustic, PIR", "ไม่ต้องสวมใส่", "ติดตั้งเฉพาะจุด, noise, blind zone"),
    ("Radio/WiFi", "RSSI, CSI, radar", "passive, privacy-friendly", "domain shift, multi-person, calibration"),
]
for r, row in enumerate(rows):
    y = y0 + row_h + Inches(r * 0.73)
    for c, value in enumerate(row):
        fill = WHITE if r % 2 == 0 else rgb("#F8FAFD")
        color = INK if c == 0 else MUTED
        bold = c == 0
        table_cell(s, x0 + Inches(sum(widths[:c])), y, Inches(widths[c]), Inches(0.73), value, fill, color, 9.8, bold)
rounded(s, Inches(0.95), Inches(5.78), Inches(11.35), Inches(0.67), rgb("#EAF4FF"), None)
center_text(s, "ประเด็นจาก review: accuracy ใน lab ไม่พอ ต้องมี real-world testing และ unobtrusive UX", Inches(1.22), Inches(5.96), Inches(10.8), Inches(0.24), 14.5, BLUE, True)
source(s, "Chaudhuri et al., 2014; Petersen et al., 2020; fall-detection scoping reviews")
footer(s, "3:30 - 4:50")
add_note(5, "ภาพรวม existing solutions", "3:30 - 4:50", "อธิบาย taxonomy ไม่ต้องลง paper ทุกตัว เน้น trade-off ว่าระบบที่ใช้งานจริงต้อง balance privacy, compliance, false alarm และ coverage")


# 6 Limitations matrix
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "ช่องว่าง", 6)
headline(s, "ช่องว่างของวิธีเดิมอยู่ที่การใช้งานจริง", "ระบบที่ดีบนกระดาษอาจพังเพราะ privacy, การลืมใส่, หรือการต้องกดเอง")
cards = [
    ("Camera", "ตรวจจับ posture ได้ดี\nแต่พื้นที่ส่วนตัว เช่น ห้องนอน/ห้องน้ำ ทำให้ accept ยาก", RED),
    ("Wearable", "เหมาะกับ outdoor และ GPS\nแต่ fall detection fail ได้ถ้าไม่ได้ใส่หรือแบตหมด", BLUE),
    ("SOS button", "ง่ายและต้นทุนต่ำ\nแต่ต้องมีสติและมีแรงกดปุ่ม", AMBER),
    ("Dense ambient sensors", "ไม่ต้องสวมใส่\nแต่ติดตั้งเยอะและครอบคลุมพื้นที่ยาก", PURPLE),
]
for i, (title, body, color) in enumerate(cards):
    x = Inches(0.86 + (i % 2) * 6.0)
    y = Inches(2.0 + (i // 2) * 1.72)
    card(s, x, y, Inches(5.45), Inches(1.25), title, body, color, 15, 11.5)
rounded(s, Inches(0.98), Inches(5.75), Inches(11.25), Inches(0.78), BLACK, None)
center_text(s, "Design target: ไม่บันทึกภาพ, ไม่ต้องสวมใส่, ไม่ต้องให้ผู้สูงอายุเป็นคนเริ่ม action", Inches(1.25), Inches(5.95), Inches(10.72), Inches(0.28), 15.2, WHITE, True)
source(s, "Chaudhuri et al., 2014; wearable umbrella review, 2020")
footer(s, "4:50 - 6:10")
add_note(6, "ข้อจำกัดของวิธีเดิม", "4:50 - 6:10", "โยงว่า pain point ไม่ใช่แค่ model accuracy แต่เป็น user adoption และ emergency workflow")


# 7 Why CSI
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "แนวคิดหลัก", 7)
headline(s, "เราจึงเลือก WiFi CSI เป็นตัวตรวจจับที่ไม่ใช้กล้อง", "CSI เห็นการเปลี่ยนแปลงของ channel หลาย subcarrier เมื่อตัวคนรบกวน multipath")
text(s, "Camera sees appearance\nWearable sees body motion\nWiFi CSI sees channel disturbance", Inches(0.86), Inches(2.0), Inches(4.7), Inches(1.25), 24, INK, True, font=LATIN)
rect(s, Inches(0.9), Inches(3.55), Inches(3.7), Inches(0.08), RED)
bullets(s, [
    "Passive: คนไม่ต้องถือหรือใส่อุปกรณ์",
    "Privacy: ไม่มีภาพหน้า/ร่างกาย",
    "Low-cost: ESP32 และ WiFi hardware",
    "Works in low light: ไม่พึ่งแสง",
], Inches(0.9), Inches(4.05), Inches(5.1), 13, 0.45, MUTED, BLUE)
rounded(s, Inches(7.0), Inches(2.02), Inches(5.25), Inches(3.9), WHITE, LINE)
center_text(s, "Tx", Inches(7.35), Inches(2.55), Inches(0.7), Inches(0.28), 15, WHITE, True, LATIN)
circle(s, Inches(7.28), Inches(2.28), Inches(0.86), BLUE, None)
circle(s, Inches(10.95), Inches(4.62), Inches(0.86), TEAL, None)
center_text(s, "Rx", Inches(11.02), Inches(4.89), Inches(0.7), Inches(0.28), 15, WHITE, True, LATIN)
for offset, color in [(0.0, RED), (0.35, AMBER), (0.7, PURPLE), (1.05, GREEN)]:
    line(s, Inches(8.05), Inches(2.72 + offset), Inches(10.95), Inches(4.97 - offset / 2), color, 2)
rounded(s, Inches(9.25), Inches(3.05), Inches(1.0), Inches(1.35), rgb("#FFE4D6"), None)
center_text(s, "คน\nเคลื่อนไหว", Inches(9.25), Inches(3.34), Inches(1.0), Inches(0.52), 12.5, INK, True)
source(s, "Wi-ESP, 2020; Espressif ESP-CSI docs; ESP32-CSI-Tool")
footer(s, "6:10 - 7:20")
add_note(7, "เหตุผลที่เลือก CSI", "6:10 - 7:20", "อธิบายแบบ intuition: ไม่ได้เห็นตัวคน แต่เห็น channel ที่เปลี่ยนเมื่อร่างกายทำให้สัญญาณ WiFi สะท้อนและดูดกลืนต่างไป")


# 8 CSI pipeline
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "สัญญาณและข้อมูล", 8)
headline(s, "สัญญาณ WiFi ถูกแปลงเป็นคุณลักษณะก่อนเข้าโมเดล", "pipeline ของโปรเจคนี้ใช้ 52 subcarriers และสถิติต่อ subcarrier")
steps = [
    ("ESP32", "CSI packets\n100 pkt/s", RED),
    ("Window", "200 packets\nsliding buffer", BLUE),
    ("Feature", "8 stats × 52\n= 416", TEAL),
    ("Model", "fall / no_fall\nconfidence", GREEN),
    ("Alert", "push + ack\nor call", AMBER),
]
for i, (label, body, color) in enumerate(steps):
    x = Inches(0.75 + i * 2.48)
    flow_step(s, x, Inches(2.35), label, body, color, 1.8)
    if i < len(steps) - 1:
        line(s, x + Inches(1.55), Inches(2.86), x + Inches(2.18), Inches(2.86), rgb("#BAC5D2"), 3)
rounded(s, Inches(1.15), Inches(5.15), Inches(10.95), Inches(0.85), WHITE, LINE)
center_text(s, "Feature set: mean, min, max, std, MAD, IQR, skewness, kurtosis ต่อ subcarrier", Inches(1.4), Inches(5.37), Inches(10.45), Inches(0.28), 15, INK, True)
source(s, "Local source: fall_detection_backend/ml_service/app/preprocess.py")
footer(s, "7:20 - 8:30")
add_note(8, "Pipeline ของสัญญาณ", "7:20 - 8:30", "เน้นว่าเราไม่ส่ง raw ทุก packet เข้า model ตรง ๆ แต่สรุปเป็น window features เพื่อให้ inference เร็วและเบากว่า sequence หนัก ๆ")


# 9 Related work WiFall/RT-Fall
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "งานวิจัยที่เกี่ยวข้อง", 9)
headline(s, "งานวิจัยยืนยันว่า WiFi CSI ตรวจจับการล้มได้", "แต่หลายงานยังผูกกับ environment และ hardware setup ที่ควบคุมไว้")
paper_cards = [
    ("WiFall", "ใช้ CSI เป็น indicator ของ activity\nชี้ว่าตรวจจับ fall แบบไม่ใส่อุปกรณ์เป็นไปได้", BLUE),
    ("Anti-Fall", "ใช้ CSI phase/amplitude\nเน้น real-time และแยก fall-like activities", RED),
    ("RT-Fall", "commodity WiFi, contactless\nรายงาน sensitivity/specificity สูงกว่า WiFall ในหลาย scenario", TEAL),
]
for i, (title, body, color) in enumerate(paper_cards):
    card(s, Inches(0.8 + i * 4.15), Inches(2.05), Inches(3.55), Inches(2.0), title, body, color, 16, 12)
rounded(s, Inches(1.05), Inches(4.95), Inches(11.25), Inches(0.9), BLACK, None)
center_text(s, "บทเรียน: CSI มี signal สำหรับ fall จริง แต่ต้องระวัง overfit กับห้อง, มุมวาง, คนทดสอบ และกิจกรรมคล้ายล้ม", Inches(1.32), Inches(5.18), Inches(10.7), Inches(0.3), 15, WHITE, True)
source(s, "WiFall; Anti-Fall, 2015; RT-Fall, IEEE TMC 2017")
footer(s, "8:30 - 10:00")
add_note(9, "Related work fall detection", "8:30 - 10:00", "สรุปว่ามี precedent เชิงวิชาการ แต่ต้องไม่ claim ว่าทุก environment จะใช้ได้ทันที เพราะ domain shift เป็นปัญหาหลัก")


# 10 Related work ML + ESP32
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "งานวิจัยที่เกี่ยวข้อง", 10)
headline(s, "งานฝั่ง ESP32 ทำให้แนวคิดนี้ต้นทุนต่ำขึ้น", "นี่ตรงกับโจทย์ของเรา: low-cost prototype ที่เอาไปติดตั้งได้จริงมากขึ้น")
card(s, Inches(0.82), Inches(2.0), Inches(3.7), Inches(2.3), "Device-free HAR + LSTM", "งาน CSI-HAR ใช้ SVM/LSTM เพื่อจำแนก activity เช่น walk, sit, stand และ fall\nแสดงว่า temporal model มีประโยชน์", BLUE, 14, 11)
card(s, Inches(4.83), Inches(2.0), Inches(3.7), Inches(2.3), "ESP32 CSI Tool", "ESP32 เก็บ CSI จาก WiFi-enabled microcontroller ได้\nเหมาะกับ prototype ราคาต่ำ", TEAL, 14, 11)
card(s, Inches(8.84), Inches(2.0), Inches(3.7), Inches(2.3), "Wi-ESP / ESP-CSI", " literature และ ecosystem เริ่มรองรับ ESP32 สำหรับ device-free WiFi sensing มากขึ้น", GREEN, 14, 11)
bullets(s, [
    "จุดที่เราใช้ตาม literature: CSI windowing, feature extraction, binary fall decision",
    "จุดที่ต้องพิสูจน์เอง: reliability ในห้องจริงและ unseen scenario",
], Inches(1.0), Inches(5.08), Inches(11.4), 13.5, 0.5, INK, RED)
source(s, "Damodaran et al., 2020; ESP32-CSI-Tool; Wi-ESP, 2020; Espressif ESP-CSI")
footer(s, "10:00 - 11:20")
add_note(10, "Related work ESP32/ML", "10:00 - 11:20", "โยงว่างานของเราไม่ได้เริ่มจากศูนย์ แต่เลือกส่วนที่พิสูจน์แล้วมาทำ prototype end-to-end")


# 11 Gap and contribution
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "สิ่งที่งานนี้เติม", 11)
headline(s, "สิ่งที่เราเติมคือระบบแจ้งเตือนครบวงจร", "ไม่ใช่แค่ classifier แต่เป็น safety workflow ที่ตรวจจับแล้วต้องมีคนตอบรับ")
gaps = [
    ("Lab accuracy", "หลาย paper รายงาน accuracy สูง แต่ยังต้องทดสอบ cross-room / cross-person"),
    ("Alert workflow", "งาน ML จำนวนมากจบที่ prediction แต่ระบบจริงต้องมี acknowledge และ escalation"),
    ("Low-cost deployment", "ESP32 ทำให้ต้นทุนต่ำ แต่ signal quality และ calibration ยังเป็นความเสี่ยง"),
]
for i, (title, body) in enumerate(gaps):
    card(s, Inches(0.9 + i * 4.05), Inches(2.0), Inches(3.45), Inches(2.05), title, body, [RED, BLUE, TEAL][i], 14, 11.4)
rounded(s, Inches(1.15), Inches(5.1), Inches(10.95), Inches(0.85), rgb("#FFE9EE"), None)
center_text(s, "Contribution ของเรา: prototype end-to-end จาก ESP32 CSI → ML inference → app notification → caregiver ack → Twilio escalation", Inches(1.42), Inches(5.33), Inches(10.4), Inches(0.28), 15, RED_DARK, True)
footer(s, "11:20 - 12:20")
add_note(11, "Gap และ contribution", "11:20 - 12:20", "บอกให้ชัดว่า contribution ไม่ได้อ้างว่าแก้ปัญหา CSI ทั้งโลก แต่ทำระบบ end-to-end ที่เห็น path ไปสู่ deployment")


# 12 Solution overview
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "แนวทางของเรา", 12)
headline(s, "แนวทางของเราคือตรวจจับแล้วส่งถึงผู้ดูแลทันที", "ตรวจจับด้วย WiFi CSI และส่งต่อไปยังแอปของผู้ดูแล")
steps = [
    ("SENSE", "ESP32 รับ CSI", RED),
    ("INFER", "ML service ทำนาย", BLUE),
    ("STORE", "Supabase log", PURPLE),
    ("NOTIFY", "Expo push", TEAL),
    ("ESCALATE", "Twilio call", AMBER),
]
for i, (label, body, color) in enumerate(steps):
    x = Inches(0.75 + i * 2.48)
    circle(s, x + Inches(0.4), Inches(2.28), Inches(1.08), color, None)
    center_text(s, label, x + Inches(0.4), Inches(2.68), Inches(1.08), Inches(0.18), 8.5, WHITE, True, LATIN)
    center_text(s, body, x, Inches(3.6), Inches(1.9), Inches(0.38), 12.5, INK, True)
    if i < 4:
        line(s, x + Inches(1.55), Inches(2.82), x + Inches(2.17), Inches(2.82), rgb("#B9C4D1"), 3)
rounded(s, Inches(1.35), Inches(5.0), Inches(4.85), Inches(0.8), rgb("#EAF4FF"), None)
center_text(s, "Ack ทันเวลา → จบ flow", Inches(1.55), Inches(5.23), Inches(4.45), Inches(0.25), 16, BLUE, True)
rounded(s, Inches(7.0), Inches(5.0), Inches(4.85), Inches(0.8), rgb("#FFE8EE"), None)
center_text(s, "ไม่ Ack → โทรฉุกเฉิน", Inches(7.2), Inches(5.23), Inches(4.45), Inches(0.25), 16, RED_DARK, True)
footer(s, "12:20 - 13:20")
add_note(12, "Solution overview", "12:20 - 13:20", "เล่า flow หนึ่งรอบจากสัญญาณจนถึงการช่วยเหลือจริง เพื่อให้คนฟังเห็นระบบก่อนดู implementation")


# 13 Architecture
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "สถาปัตยกรรม", 13)
headline(s, "ระบบปัจจุบันเชื่อมครบตั้งแต่ ESP32 ถึงแอป", "แบ่งเป็น hardware, ML service, Express API, database, และ Expo mobile app")
layers = [
    ("Hardware", "ESP32 AP/STA\nCSI via USB serial", RED),
    ("ML service", "FastAPI\nRF active inference", BLUE),
    ("Express API", "predict/events/push/alert\nAPI key auth", PURPLE),
    ("Supabase", "fall_events\npush_tokens", GREEN),
    ("Mobile", "Expo app\npoll + ack + push", TEAL),
]
for i, (title, body, color) in enumerate(layers):
    x = Inches(0.62 + i * 2.54)
    card(s, x, Inches(2.1), Inches(2.18), Inches(1.7), title, body, color, 12.5, 9.5)
    if i < len(layers) - 1:
        line(s, x + Inches(2.18), Inches(2.95), x + Inches(2.42), Inches(2.95), rgb("#B8C3CF"), 2.4)
rounded(s, Inches(0.82), Inches(4.65), Inches(11.7), Inches(1.25), WHITE, LINE)
text(s, "Key endpoints", Inches(1.08), Inches(4.9), Inches(2.0), Inches(0.3), 13.5, INK, True, font=LATIN)
text(s, "POST /api/v1/predict   GET /api/v1/events/falls   POST /api/v1/alert/ack/:event_id   POST /api/v1/push/register   POST /api/v1/demo/fire", Inches(2.75), Inches(4.9), Inches(9.35), Inches(0.45), 10.3, MUTED, False, font=LATIN)
source(s, "Local source: express_api/src/index.js, routes/*.js, lib/api.ts")
footer(s, "13:20 - 14:40")
add_note(13, "Architecture", "13:20 - 14:40", "ย้ำว่าใช้ source code ปัจจุบัน: ML service active เป็น RF, API มี demo endpoint, app poll event และกด ack ได้")


# 14 Data feature pipeline
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "ข้อมูล", 14)
headline(s, "เราเก็บ CSI ผ่าน USB serial เพื่อให้ข้อมูลนิ่งขึ้น", "โหมดปัจจุบันเน้น USB serial เพื่อให้เก็บ CSI เสถียรกว่า UDP/Hotspot")
card(s, Inches(0.85), Inches(2.0), Inches(3.55), Inches(2.2), "Collection", "ESP32 receiver ส่ง CSI ผ่าน USB serial\nbaud 921600\ncollector: csi_collector_serial.py", RED, 14, 11)
card(s, Inches(4.88), Inches(2.0), Inches(3.55), Inches(2.2), "Preprocess", "parse I/Q หรือ amplitude\n52 subcarriers\nwindow 200 packets", BLUE, 14, 11)
card(s, Inches(8.9), Inches(2.0), Inches(3.55), Inches(2.2), "Feature/Train", "416 features ต่อ window\ntrain notebooks + model files\nRF/LSTM artifacts", TEAL, 14, 11)
rounded(s, Inches(1.05), Inches(5.05), Inches(11.25), Inches(0.85), WHITE, LINE)
center_text(s, "Decision ที่เปลี่ยนไป: จาก sequence-heavy LSTM มาใช้ RF window ล่าสุด เพื่อให้ real-time ง่ายขึ้นและลด buffer issue", Inches(1.3), Inches(5.28), Inches(10.75), Inches(0.28), 14.6, INK, True)
source(s, "Local source: preprocess.py, predict.py, data_collection/csi_collector_serial.py")
footer(s, "14:40 - 15:50")
add_note(14, "Data pipeline", "14:40 - 15:50", "พูดเรื่องการตัดสินใจทางวิศวกรรม: USB serial ลด dependency ต่อ hotspot และ RF ใช้ window ล่าสุดเพื่อ response time")


# 15 Model journey
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "โมเดล", 15)
headline(s, "ผลทดสอบกับชุดข้อมูลเดิมดี แต่ยังต้องพิสูจน์ในสถานการณ์จริง", "นี่คือบทเรียนสำคัญจาก lit review และจากงานเราเอง")
labels = [
    ("v1", "LSTM binary\n87.3%", RED),
    ("v2", "LSTM 4-class\n92.9%", BLUE),
    ("v3", "LSTM binary\n99.7% split", TEAL),
    ("active", "RF v1\nthreshold 0.55", GREEN),
]
for i, (title, body, color) in enumerate(labels):
    x = Inches(0.95 + i * 2.95)
    rounded(s, x, Inches(2.25), Inches(2.35), Inches(1.55), WHITE, LINE)
    pill(s, title, x + Inches(0.22), Inches(2.47), color, Inches(0.9))
    center_text(s, body, x + Inches(0.18), Inches(3.0), Inches(1.98), Inches(0.48), 13, INK, True)
    if i < 3:
        line(s, x + Inches(2.35), Inches(3.03), x + Inches(2.75), Inches(3.03), rgb("#B8C3CF"), 2.3)
rounded(s, Inches(1.08), Inches(4.85), Inches(11.15), Inches(0.95), rgb("#FFF4E5"), None)
center_text(s, "ข้อสรุป: offline accuracy สูงเป็นหลักฐานว่า signal มี pattern แต่ยังไม่ใช่ evidence ว่าใช้ได้ทุกบ้าน", Inches(1.35), Inches(5.1), Inches(10.6), Inches(0.3), 15.5, AMBER, True)
source(s, "Local model reports: training_report.json, report_v2.json, report_v3.json; active source: predict.py")
footer(s, "15:50 - 17:00")
add_note(15, "Model journey", "15:50 - 17:00", "พูดตรง ๆ ว่าตัวเลข test split ดี แต่เราเรียนรู้ว่า real-time และ unseen environment สำคัญกว่า จึงต้องมี evaluation plan ใหม่")


# 16 Active inference
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "โมเดลที่ใช้งาน", 16)
headline(s, "โมเดลที่ใช้งานตอนนี้คือ Random Forest จากหน้าต่างข้อมูลล่าสุด", "ML service รับ sequence list แต่ใช้ window ล่าสุดเพื่อ classify")
card(s, Inches(0.9), Inches(2.0), Inches(3.55), Inches(2.35), "Input", "features: List[List[float]]\nใช้ features[-1]\nshape 1 × 416", BLUE, 14, 11.5)
card(s, Inches(4.86), Inches(2.0), Inches(3.55), Inches(2.35), "Model", "rf_v1.pkl + scaler_rf_v1.pkl\npredict_proba\nfall_prob = proba[0]", TEAL, 14, 11.5)
card(s, Inches(8.82), Inches(2.0), Inches(3.55), Inches(2.35), "Decision", "fall_prob >= 0.55\nreturn fall/no_fall\nconfidence score", RED, 14, 11.5)
bullets(s, [
    "ข้อดี: latency ต่ำ, implementation ง่าย, เหมาะกับ demo และ real-time loop",
    "ข้อควรระวัง: threshold ยังต้อง tune ด้วย ROC/PR curve และ unseen test",
], Inches(1.05), Inches(5.05), Inches(11.0), 13.5, 0.48, INK, RED)
source(s, "Local source: fall_detection_backend/ml_service/app/predict.py")
footer(s, "17:00 - 18:00")
add_note(16, "Active inference", "17:00 - 18:00", "ย้ำว่า slide นี้คือ source of truth ปัจจุบัน ไม่ใช่ README เก่าที่พูดว่า LSTM active")


# 17 Alert flow
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "ระบบแจ้งเตือน", 17)
headline(s, "หลังตรวจพบการล้ม ระบบรอผู้ดูแลกดรับทราบ", "ถ้าไม่มีการตอบรับในเวลาที่กำหนด ระบบจะยกระดับไปยังการโทรฉุกเฉิน")
states = [
    ("DETECTED", "save DB\nsend push", RED),
    ("AWAITING_ACK", "timer 60s\ncooldown", BLUE),
    ("ACKED", "cancel timer\nmark ack", GREEN),
    ("ESCALATED", "Twilio voice\nmark escalated", AMBER),
]
for i, (label, body, color) in enumerate(states):
    x = Inches(0.85 + i * 3.05)
    rounded(s, x, Inches(2.25), Inches(2.42), Inches(1.35), WHITE, LINE)
    pill(s, label, x + Inches(0.22), Inches(2.47), color, Inches(1.65))
    center_text(s, body, x + Inches(0.18), Inches(3.02), Inches(2.06), Inches(0.38), 11.2, INK, True)
    if i < 1:
        line(s, x + Inches(2.42), Inches(2.92), x + Inches(2.85), Inches(2.92), rgb("#B8C3CF"), 2.3)
line(s, Inches(6.28), Inches(3.68), Inches(7.0), Inches(4.45), GREEN, 2.6)
line(s, Inches(6.28), Inches(3.68), Inches(10.05), Inches(4.45), AMBER, 2.6)
center_text(s, "กดรับทราบ", Inches(6.6), Inches(4.45), Inches(1.45), Inches(0.25), 10.5, GREEN, True)
center_text(s, "timeout", Inches(9.45), Inches(4.45), Inches(1.2), Inches(0.25), 10.5, AMBER, True, LATIN)
rounded(s, Inches(1.12), Inches(5.55), Inches(10.9), Inches(0.55), rgb("#EAF4FF"), None)
center_text(s, "ข้อดีของ binary + acknowledge: ลด false escalation และให้ caregiver เป็นคนตัดสินบริบทจริง", Inches(1.4), Inches(5.72), Inches(10.35), Inches(0.2), 13.8, BLUE, True)
source(s, "Local source: routes/predict.js, routes/alert.js, services/escalationService.js")
footer(s, "18:00 - 19:10")
add_note(17, "Alert workflow", "18:00 - 19:10", "อธิบายว่าระบบ safety ต้องมี feedback loop: detect, notify, acknowledge, escalate ไม่ใช่แค่มี model")


# 18 Mobile status
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "แอปผู้ดูแล", 18)
headline(s, "แอป Middle พร้อมเป็นหน้าจอหลักของผู้ดูแล", "มี alert list, notification, acknowledge และ screen flow สำหรับบ้าน/อุปกรณ์")
card(s, Inches(0.86), Inches(2.0), Inches(3.65), Inches(2.15), "Live alerts", "poll /events/falls ทุก 4 วินาที\nmerge เข้า AlertsContext\nแสดง active/completed/no_response", RED, 14, 11)
card(s, Inches(4.84), Inches(2.0), Inches(3.65), Inches(2.15), "Acknowledge", "caregiver กด confirm\nPOST /alert/ack/:event_id\nยกเลิก escalation timer", GREEN, 14, 11)
card(s, Inches(8.82), Inches(2.0), Inches(3.65), Inches(2.15), "Push", "ลงทะเบียน Expo token\nรับ notification tap\nroute กลับ /home", BLUE, 14, 11)
bullets(s, [
    "UI screens พร้อมสำหรับบ้าน, members, devices, contacts, notifications",
    "ข้อค้าง: API URL/API key และ Supabase config ยังเป็น development defaults",
], Inches(1.0), Inches(5.0), Inches(11.2), 13.2, 0.46, INK, RED)
source(s, "Local source: app/home.tsx, app/_layout.tsx, hooks/useLiveAlerts.ts, lib/api.ts")
footer(s, "19:10 - 20:10")
add_note(18, "Mobile progress", "19:10 - 20:10", "พูดจากมุมผู้ใช้: caregiver เห็น alert, กดรับทราบ และรับ push ได้แล้ว แต่ยังไม่ใช่ production config")


# 19 Backend status
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "ระบบหลังบ้าน", 19)
headline(s, "ระบบหลังบ้านประสานโมเดล ฐานข้อมูล push และการโทร", "Express API เป็นตัวควบคุมลำดับการแจ้งเตือนฉุกเฉิน")
left_items = [
    ("Predict route", "validate input, call ML, create event only on fall", GREEN),
    ("Events route", "list events และ fall-only events", GREEN),
    ("Push route", "register token, send test fall notification", GREEN),
    ("Alert route", "ack event, test Twilio SMS/call", GREEN),
    ("Demo route", "fire event โดยข้าม ML เพื่อโชว์ end-to-end", GREEN),
]
for i, (title, body, color) in enumerate(left_items):
    y = Inches(1.9 + i * 0.78)
    pill(s, "done", Inches(0.85), y + Inches(0.04), color, Inches(0.65))
    text(s, title, Inches(1.65), y, Inches(2.55), Inches(0.28), 12.5, INK, True, font=LATIN)
    text(s, body, Inches(4.05), y, Inches(7.7), Inches(0.3), 11.3, MUTED)
rounded(s, Inches(1.0), Inches(5.95), Inches(11.1), Inches(0.5), rgb("#FFE9EE"), None)
center_text(s, "Production risk: timers และ cooldown ยังอยู่ใน memory ถ้า server restart pending escalation จะหาย", Inches(1.25), Inches(6.08), Inches(10.6), Inches(0.18), 12.5, RED_DARK, True)
source(s, "Local source: express_api/src/routes/*.js, services/*.js")
footer(s, "20:10 - 21:10")
add_note(19, "Backend progress", "20:10 - 21:10", "สรุปว่า backend ใช้งานได้ครบสำหรับ demo แต่ต้องพูด limitation เรื่อง durability ของ timer")


# 20 Hardware status
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "อุปกรณ์", 20)
headline(s, "อุปกรณ์พร้อมโชว์สัญญาณ CSI ดิบและเก็บข้อมูลเพิ่ม", "serial mode เหมาะกับการเก็บ dataset และ demo raw signal")
card(s, Inches(0.9), Inches(2.0), Inches(3.55), Inches(2.15), "ESP32 AP/STA", "receiver สร้าง CSI network\nsender ส่ง packet ต่อเนื่อง\nตั้ง channel ให้คุม environment", RED, 14, 11)
card(s, Inches(4.86), Inches(2.0), Inches(3.55), Inches(2.15), "USB Serial", "อ่าน CSI ที่ 921600 baud\nลดปัญหา firewall/hotspot\nMac ยังต่อ WiFi ปกติได้", BLUE, 14, 11)
card(s, Inches(8.82), Inches(2.0), Inches(3.55), Inches(2.15), "Live demo", "โชว์ packet flow + raw signal response\nแยกชัดว่า raw signal ไม่ใช่ proof ของ model", TEAL, 14, 11)
rounded(s, Inches(1.08), Inches(5.02), Inches(11.15), Inches(0.82), WHITE, LINE)
center_text(s, "สิ่งที่ทำแล้ว: collector, ESP checker, preprocess pipeline, model artifacts และคู่มือ live demo", Inches(1.35), Inches(5.26), Inches(10.6), Inches(0.27), 15, INK, True)
source(s, "Local source: CLAUDE.md, csi_collector_serial.py, esp_checker.py, live_demo_manual.md")
footer(s, "21:10 - 22:10")
add_note(20, "Hardware progress", "21:10 - 22:10", "ถ้าจะ demo ให้บอกตรง ๆ ว่า live hardware demo คือโชว์ raw CSI response ไม่ใช่ model generalization")


# 21 Done summary
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "ความคืบหน้า", 21)
headline(s, "ตอนนี้เรามีต้นแบบครบวงจรแล้ว", "มองเป็น deliverable แทนแค่รายการไฟล์")
done = [
    ("ML", "preprocess + RF active inference + LSTM/RF artifacts"),
    ("Backend", "predict, events, ack, push, demo, Twilio integration"),
    ("Database", "Supabase schema: fall_events, push_tokens"),
    ("Mobile", "alert list, acknowledge, notification registration"),
    ("Hardware", "ESP32 CSI collection through USB serial"),
    ("Demo", "manual + shortcut endpoint for end-to-end presentation"),
]
for i, (title, body) in enumerate(done):
    x = Inches(0.82 + (i % 3) * 4.0)
    y = Inches(1.95 + (i // 3) * 1.52)
    card(s, x, y, Inches(3.48), Inches(1.13), title, body, [RED, BLUE, TEAL, GREEN, PURPLE, AMBER][i], 13.5, 10.5)
rounded(s, Inches(1.1), Inches(5.65), Inches(11.15), Inches(0.65), BLACK, None)
center_text(s, "สถานะจริง: prototype end-to-end พร้อมเล่าเป็นระบบ แต่ evaluation real-world ยังเป็นงานหลักที่ต้องทำต่อ", Inches(1.35), Inches(5.84), Inches(10.65), Inches(0.22), 14.2, WHITE, True)
footer(s, "22:10 - 23:00")
add_note(21, "Summary done", "22:10 - 23:00", "ใช้ slide นี้เป็น milestone summary ก่อนเข้าสู่ limitation เพื่อไม่ให้ presentation ดูเหมือนมีแต่ปัญหา")


# 22 Limitations
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "ข้อจำกัด", 22)
headline(s, "ข้อจำกัดสำคัญคือความแม่นยำเมื่อเปลี่ยนสภาพแวดล้อม", "ต้องแยกผล offline, demo และ real-world validation ออกจากกันให้ชัด")
limits = [
    ("Domain shift", "โมเดลอาจยึดติดกับห้อง มุมวาง และคนที่ใช้ train"),
    ("False alarms", "กิจกรรมคล้ายล้ม เช่น นั่งลงเร็วหรือก้ม อาจ trigger"),
    ("Durability", "ack timer ยังเป็น in-memory ต้องใช้ Redis/BullMQ สำหรับ production"),
    ("Configuration", "API key, API base URL, Supabase config ยังเป็น dev defaults"),
]
for i, (title, body) in enumerate(limits):
    x = Inches(0.88 + (i % 2) * 6.0)
    y = Inches(1.95 + (i // 2) * 1.65)
    card(s, x, y, Inches(5.45), Inches(1.18), title, body, [RED, AMBER, BLUE, PURPLE][i], 14, 11.2)
rounded(s, Inches(1.1), Inches(5.6), Inches(11.15), Inches(0.72), rgb("#EAF4FF"), None)
center_text(s, "คำตอบเชิงวิศวกรรม: แยก demo, offline metric, และ real-world validation ออกจากกันให้ชัด", Inches(1.35), Inches(5.8), Inches(10.65), Inches(0.25), 14.8, BLUE, True)
footer(s, "23:00 - 23:40")
add_note(22, "Limitations", "23:00 - 23:40", "พูด limitation ก่อนกรรมการถาม เพื่อแสดงว่าเข้าใจทั้ง ML risk และ system risk")


# 23 Next steps
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "แผนต่อไป", 23)
headline(s, "งานต่อไปคือพิสูจน์ความน่าเชื่อถือก่อนใช้งานจริง", "เอา gap จาก literature มาแปลงเป็น evaluation plan")
steps = [
    ("1", "Unseen test", "แยกคน ห้อง มุมวาง และท่าคล้ายล้มออกจาก train", RED),
    ("2", "Threshold tuning", "ROC/PR curve เพื่อเลือก threshold จาก false alarm target", BLUE),
    ("3", "Real-time trial", "วัด latency, missed fall, false alarm ใน session จริง", TEAL),
    ("4", "Production hardening", "persistent queue/timer, env config, deployment", GREEN),
]
for i, (num, title, body, color) in enumerate(steps):
    x = Inches(0.82 + i * 3.05)
    rounded(s, x, Inches(2.05), Inches(2.48), Inches(2.45), WHITE, LINE)
    circle(s, x + Inches(0.72), Inches(2.35), Inches(0.78), color, None)
    center_text(s, num, x + Inches(0.72), Inches(2.58), Inches(0.78), Inches(0.16), 11, WHITE, True, LATIN)
    center_text(s, title, x + Inches(0.22), Inches(3.35), Inches(2.05), Inches(0.28), 13.5, INK, True)
    center_text(s, body, x + Inches(0.22), Inches(3.75), Inches(2.05), Inches(0.48), 10.3, MUTED)
rounded(s, Inches(1.22), Inches(5.55), Inches(10.9), Inches(0.62), BLACK, None)
center_text(s, "Success criteria: Fall recall สูง, false alarm ต่ำ, และ caregiver workflow ทำงานครบในสถานการณ์จริง", Inches(1.48), Inches(5.73), Inches(10.38), Inches(0.22), 14.2, WHITE, True)
footer(s, "23:40 - 24:00")
add_note(23, "Next steps", "23:40 - 24:00", "ปิดเนื้อหาด้วยแผนที่ตอบ gap โดยตรง แล้วเตรียมเข้าสู่ Q&A")


# 24 Q&A
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLACK)
kicker(s, "ถามตอบ", 24)
text(s, "พร้อมตอบคำถาม", Inches(0.8), Inches(1.0), Inches(5.5), Inches(0.8), 38, WHITE, True)
text(s, "ประเด็นที่เตรียมตอบ", Inches(0.84), Inches(2.15), Inches(3.2), Inches(0.35), 17, rgb("#E5EDF5"), True)
bullets(s, [
    "ทำไมเลือก WiFi CSI แทนกล้องหรือ wearable",
    "ตัวเลข accuracy กับ real-world reliability ต่างกันอย่างไร",
    "ตอนนี้ active model เป็นอะไร และทำไมเปลี่ยนจาก LSTM",
    "ถ้า false alarm หรือ server restart ระบบจะทำอย่างไรต่อ",
], Inches(0.9), Inches(2.75), Inches(6.2), 14, 0.55, rgb("#D6E0EA"), RED)
rounded(s, Inches(8.1), Inches(1.6), Inches(3.7), Inches(3.9), rgb("#F4F7FB"), None)
center_text(s, "24:00 - 30:00", Inches(8.5), Inches(2.42), Inches(2.9), Inches(0.45), 24, RED, True, LATIN)
center_text(s, "Q&A\n+ backup references", Inches(8.5), Inches(3.15), Inches(2.9), Inches(0.8), 24, INK, True)
footer(s, "24:00 - 30:00")
add_note(24, "Q&A", "24:00 - 30:00", "เปิดให้ถามตอบ 6 นาที ถ้าถาม references หรือ implementation details ให้เปิด slide ถัดไปหรือชี้ไปที่ repo/source")


# 25 References
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "อ้างอิง", 25)
headline(s, "แหล่งอ้างอิงหลัก", "สไลด์นี้เป็น backup สำหรับคำถามเรื่องที่มาและ literature")
refs = [
    "WHO. Falls fact sheet. 2021. https://www.who.int/en/news-room/fact-sheets/detail/falls",
    "CDC. Facts About Falls. 2026. https://www.cdc.gov/falls/data-research/facts-stats/",
    "Chaudhuri, S., Thompson, H., & Demiris, G. Fall Detection Devices and their Use with Older Adults. 2014.",
    "Petersen et al. Are wearable devices effective for preventing and detecting falls: umbrella review. 2020.",
    "Wang et al. RT-Fall: A Real-Time and Contactless Fall Detection System with Commodity WiFi Devices. IEEE TMC, 2017.",
    "Zhang et al. Anti-Fall: Non-intrusive and Real-time Fall Detector Leveraging CSI. 2015.",
    "Damodaran et al. Device free human activity and fall recognition using WiFi CSI. 2020.",
    "Hernandez. ESP32-CSI-Tool. https://github.com/StevenMHernandez/ESP32-CSI-Tool",
    "Espressif. esp-csi. https://github.com/espressif/esp-csi",
    "Local repo sources: ml_service/app/predict.py, preprocess.py, express_api/src/routes, hooks/useLiveAlerts.ts",
]
for i, ref in enumerate(refs):
    y = Inches(1.82 + i * 0.45)
    rect(s, Inches(0.78), y + Inches(0.08), Inches(0.06), Inches(0.16), RED if i < 2 else BLUE)
    text(s, ref, Inches(0.95), y, Inches(11.8), Inches(0.32), 8.8, INK if i < 2 else MUTED, False, font=LATIN if ref.startswith(("WHO", "CDC", "Wang", "Zhang", "Damodaran", "Hernandez", "Espressif", "Local")) else FONT)
footer(s, "Backup")
add_note(25, "References", "Backup", "ใช้เป็นสไลด์สำรองเมื่อถูกถามที่มาของตัวเลขหรือ related work")


prs.save(OUT)

notes_path = Path(NOTES_OUT)
with notes_path.open("w", encoding="utf-8") as f:
    f.write("# TDG Lit Review + Progress Presentation Notes\n\n")
    f.write("เวลาเป้าหมาย: พรี 24 นาที + ถามตอบ 6 นาที\n\n")
    for n, title, timing, body in notes:
        f.write(f"## Slide {n}: {title}\n")
        f.write(f"- เวลา: {timing}\n")
        f.write(f"- พูดหลัก: {body}\n\n")

print(f"Saved {OUT}")
print(f"Saved {NOTES_OUT}")
