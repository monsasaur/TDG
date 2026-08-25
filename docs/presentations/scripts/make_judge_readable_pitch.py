import os
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "fall_detection_judge_readable_pitch.pptx")

W = Inches(13.333)
H = Inches(7.5)

BLACK = RGBColor(15, 18, 22)
INK = RGBColor(28, 34, 42)
MUTED = RGBColor(92, 103, 115)
WHITE = RGBColor(255, 255, 255)
BG = RGBColor(247, 249, 252)
LINE = RGBColor(222, 229, 237)
RED = RGBColor(255, 48, 85)
RED_DARK = RGBColor(194, 25, 52)
BLUE = RGBColor(35, 105, 220)
TEAL = RGBColor(0, 154, 142)
GREEN = RGBColor(28, 156, 92)
AMBER = RGBColor(240, 156, 28)
PURPLE = RGBColor(112, 89, 204)

FONT = "Thonburi"
LATIN = "Aptos"

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def rgb(hex_value):
    hex_value = hex_value.replace("#", "")
    return RGBColor(int(hex_value[:2], 16), int(hex_value[2:4], 16), int(hex_value[4:], 16))


def rect(slide, x, y, w, h, fill=WHITE, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def rounded(slide, x, y, w, h, fill=WHITE, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def circle(slide, x, y, d, fill=WHITE, line=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    return shape


def text(slide, value, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(frame, side, 0)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def center_text(slide, value, x, y, w, h, size=18, color=INK, bold=False, font=FONT):
    box = text(slide, value, x, y, w, h, size, color, bold, PP_ALIGN.CENTER, font)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def line(slide, x1, y1, x2, y2, color=LINE, width=2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    return c


def arrow(slide, x1, y1, x2, y2, color=LINE, width=2.5):
    c = line(slide, x1, y1, x2, y2, color, width)
    c.line.end_arrowhead = True
    return c


def bg(slide):
    rect(slide, 0, 0, W, H, BG)


def kicker(slide, value, n):
    text(slide, value.upper(), Inches(0.62), Inches(0.32), Inches(5), Inches(0.25), 9, RED, True, font=LATIN)
    text(slide, f"{n:02d}", Inches(12.1), Inches(0.32), Inches(0.6), Inches(0.25), 9, MUTED, True, PP_ALIGN.RIGHT, LATIN)


def footer(slide, timing):
    rect(slide, Inches(0.62), Inches(7.02), Inches(12.05), Inches(0.01), LINE)
    text(slide, timing, Inches(11.0), Inches(7.11), Inches(1.65), Inches(0.22), 8, MUTED, True, PP_ALIGN.RIGHT, LATIN)


def headline(slide, value, sub=None):
    text(slide, value, Inches(0.62), Inches(0.63), Inches(12.1), Inches(0.95), 34, INK, True)
    if sub:
        text(slide, sub, Inches(0.66), Inches(1.55), Inches(11.4), Inches(0.42), 16, MUTED)


def mini_card(slide, x, y, w, h, title, body, color):
    rounded(slide, x, y, w, h, WHITE, LINE)
    rect(slide, x, y, Inches(0.08), h, color)
    text(slide, title, x + Inches(0.23), y + Inches(0.18), w - Inches(0.42), Inches(0.32), 14, INK, True)
    text(slide, body, x + Inches(0.23), y + Inches(0.57), w - Inches(0.42), h - Inches(0.7), 11, MUTED)


def stat(slide, x, y, value, label, color):
    rounded(slide, x, y, Inches(2.25), Inches(1.15), WHITE, LINE)
    center_text(slide, value, x, y + Inches(0.16), Inches(2.25), Inches(0.38), 24, color, True, LATIN)
    center_text(slide, label, x + Inches(0.12), y + Inches(0.58), Inches(2.0), Inches(0.33), 10, MUTED, False)


# 1. Opening
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BLACK)
rect(s, Inches(7.35), 0, Inches(6.0), H, rgb("#F3F6FA"))
kicker(s, "Overview", 1)
text(s, "อัปเดตความคืบหน้าโครงการ", Inches(0.7), Inches(0.85), Inches(6.2), Inches(0.62), 33, WHITE, True)
text(s, "Fall Detection System\nด้วย WiFi CSI + AI", Inches(0.7), Inches(1.75), Inches(5.9), Inches(1.1), 31, rgb("#E8EEF5"), True)
text(s, "วันนี้จะเล่าว่า\nทำถึงไหนแล้ว", Inches(0.7), Inches(3.05), Inches(5.95), Inches(1.2), 35, WHITE, True)
rect(s, Inches(0.7), Inches(4.62), Inches(1.25), Inches(0.07), RED)
text(s, "และจะทำอะไรต่อ\nในเฟสถัดไป (Phase 2)", Inches(0.7), Inches(4.95), Inches(6.0), Inches(0.9), 22, rgb("#D9E1EA"), True)

# right visual: home, elder, alert
circle(s, Inches(8.1), Inches(1.0), Inches(3.9), WHITE, LINE)
center_text(s, "บ้าน", Inches(8.1), Inches(1.2), Inches(3.9), Inches(0.35), 18, INK, True)
rounded(s, Inches(9.12), Inches(2.0), Inches(1.85), Inches(1.35), rgb("#EAF0F7"), None)
rect(s, Inches(9.37), Inches(2.55), Inches(1.35), Inches(0.8), rgb("#D5DEE9"))
circle(s, Inches(9.65), Inches(1.58), Inches(0.75), rgb("#FFD9C7"), None)
line(s, Inches(10.02), Inches(2.35), Inches(10.02), Inches(3.05), INK, 4)
line(s, Inches(10.02), Inches(2.72), Inches(9.58), Inches(3.05), INK, 3)
line(s, Inches(10.02), Inches(2.72), Inches(10.48), Inches(3.05), INK, 3)
rounded(s, Inches(8.5), Inches(4.7), Inches(3.1), Inches(0.8), rgb("#FFE8EE"), None)
center_text(s, "FALL ALERT", Inches(8.7), Inches(4.87), Inches(2.7), Inches(0.32), 18, RED_DARK, True, LATIN)
footer(s, "0:00 - 0:35")


# 2. Problem visual
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "Recap", 2)
headline(s, "ทบทวนโจทย์เดิม", "ทุกวิธีที่มีอยู่ต้องแลกบางอย่าง: privacy, การสวมใส่, หรือการกดขอความช่วยเหลือเอง")

labels = [
    ("กล้อง", "เห็นภาพ\nแต่เสีย privacy", BLUE),
    ("นาฬิกา", "ดีถ้าใส่\nแต่ผู้สูงอายุลืมได้", TEAL),
    ("ปุ่ม SOS", "ต้องกดเอง\nถ้าหมดสติกดไม่ได้", AMBER),
]
for i, (title, body, color) in enumerate(labels):
    x = Inches(0.85 + i * 4.1)
    rounded(s, x, Inches(2.18), Inches(3.35), Inches(2.45), WHITE, LINE)
    circle(s, x + Inches(1.15), Inches(2.48), Inches(1.05), color, None)
    center_text(s, title, x + Inches(0.25), Inches(3.72), Inches(2.85), Inches(0.35), 20, INK, True)
    center_text(s, body, x + Inches(0.35), Inches(4.14), Inches(2.65), Inches(0.44), 14, MUTED)
    center_text(s, "✕", x + Inches(1.47), Inches(2.57), Inches(0.4), Inches(0.4), 28, WHITE, True, LATIN)

rounded(s, Inches(1.55), Inches(5.35), Inches(10.25), Inches(0.83), BLACK, None)
center_text(s, "โจทย์ของเรา: ตรวจจับอัตโนมัติ โดยไม่ใช้กล้อง และไม่ต้องใส่อุปกรณ์", Inches(1.8), Inches(5.55), Inches(9.75), Inches(0.36), 19, WHITE, True)
footer(s, "0:35 - 1:15")


# 3. Solution one-glance architecture
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "Current System", 3)
headline(s, "ระบบที่ทำงานอยู่ตอนนี้", "ตั้งแต่ ESP32 เก็บสัญญาณ ไปจนถึงแจ้งเตือนและโทรฉุกเฉินจริง — ทำงานครบ end-to-end แล้ว")

steps = [
    ("ESP32", "รับ WiFi CSI", RED),
    ("AI", "fall / no_fall", BLUE),
    ("APP", "แจ้งผู้ดูแล", TEAL),
    ("ACK", "กดรับทราบ", GREEN),
    ("CALL", "โทรฉุกเฉิน", AMBER),
]
for i, (name, desc, color) in enumerate(steps):
    x = Inches(0.72 + i * 2.5)
    circle(s, x + Inches(0.43), Inches(2.38), Inches(1.25), color, None)
    center_text(s, name, x + Inches(0.43), Inches(2.75), Inches(1.25), Inches(0.23), 13, WHITE, True, LATIN)
    center_text(s, desc, x, Inches(3.82), Inches(2.1), Inches(0.42), 14, INK, True)
    if i < 4:
        arrow(s, x + Inches(1.75), Inches(3.0), x + Inches(2.25), Inches(3.0), rgb("#B9C4D1"), 3)

rounded(s, Inches(0.95), Inches(5.0), Inches(5.45), Inches(0.95), rgb("#EAF4FF"), None)
center_text(s, "ถ้ากดรับทราบ → จบ flow", Inches(1.2), Inches(5.22), Inches(4.95), Inches(0.35), 18, BLUE, True)
rounded(s, Inches(6.9), Inches(5.0), Inches(5.45), Inches(0.95), rgb("#FFE8EE"), None)
center_text(s, "ถ้าไม่ตอบรับ 60 วินาที → SMS / Voice Call", Inches(7.15), Inches(5.22), Inches(4.95), Inches(0.35), 18, RED_DARK, True)
footer(s, "1:15 - 2:05")


# 4. Tech core as three proof blocks
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "Results So Far", 4)
headline(s, "ผลลัพธ์ที่ทำได้แล้ว", "มี hardware ทำงานจริง, มี model จริง, และ deploy เป็นระบบใช้งานได้จริง")

mini_card(s, Inches(0.75), Inches(2.0), Inches(3.75), Inches(2.75), "1. Signal จริง", "ESP32 รับ CSI 100 packets/sec\n52 subcarriers\nเก็บข้อมูลแล้ว 420 ไฟล์ / 1,260 windows", RED)
mini_card(s, Inches(4.82), Inches(2.0), Inches(3.75), Inches(2.75), "2. Model จริง", "LSTM v3: Accuracy 97.9%\nFall Recall 98.5%\nWindow=200, Stride=50, Seq=10", BLUE)
mini_card(s, Inches(8.9), Inches(2.0), Inches(3.75), Inches(2.75), "3. System จริง", "Backend + Supabase + Push\nAcknowledge flow\nTwilio escalation", TEAL)

rounded(s, Inches(1.05), Inches(5.55), Inches(11.25), Inches(0.72), BLACK, None)
center_text(s, "ประโยคสำคัญ: Accuracy สูงใน dataset ไม่พอ ต้องพิสูจน์กับ real-time และ unseen test", Inches(1.3), Inches(5.72), Inches(10.75), Inches(0.28), 16, WHITE, True)
footer(s, "2:05 - 3:00")


# 5. Demo storyboard
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "Demo", 5)
headline(s, "Demo ระบบที่ทำงานอยู่จริง", "โชว์เป็น 3 จังหวะ: สัญญาณนิ่ง → สัญญาณแกว่ง → แจ้งเตือนจริง")

for i, (title, desc, color) in enumerate(
    [
        ("1", "ยืนนิ่ง\nกราฟนิ่ง", TEAL),
        ("2", "เดิน\nกราฟแกว่ง", BLUE),
        ("3", "ล้ม\nกราฟกระชาก", RED),
    ]
):
    x = Inches(0.92 + i * 4.08)
    rounded(s, x, Inches(2.2), Inches(3.35), Inches(2.4), WHITE, LINE)
    circle(s, x + Inches(0.25), Inches(2.48), Inches(0.75), color, None)
    center_text(s, title, x + Inches(0.25), Inches(2.63), Inches(0.75), Inches(0.22), 15, WHITE, True, LATIN)
    center_text(s, desc, x + Inches(0.95), Inches(2.55), Inches(2.1), Inches(0.75), 21, INK, True)
    # signal line
    y = Inches(3.75)
    if title == "1":
        line(s, x + Inches(0.42), y, x + Inches(2.95), y, color, 3)
    elif title == "2":
        for k in range(6):
            line(s, x + Inches(0.42 + k * 0.38), y + Inches(0.14 if k % 2 else -0.14), x + Inches(0.64 + k * 0.38), y + Inches(-0.14 if k % 2 else 0.14), color, 3)
    else:
        line(s, x + Inches(0.42), y, x + Inches(1.35), y, color, 3)
        line(s, x + Inches(1.35), y, x + Inches(1.65), y - Inches(0.38), color, 4)
        line(s, x + Inches(1.65), y - Inches(0.38), x + Inches(2.05), y + Inches(0.35), color, 4)
        line(s, x + Inches(2.05), y + Inches(0.35), x + Inches(2.95), y + Inches(0.12), color, 3)

rounded(s, Inches(2.2), Inches(5.3), Inches(8.9), Inches(0.78), rgb("#FFE8EE"), None)
center_text(s, "จุดจบ demo: แอปขึ้น Fall Alert → ผู้ดูแลกดรับทราบ", Inches(2.5), Inches(5.48), Inches(8.3), Inches(0.3), 18, RED_DARK, True)
footer(s, "3:00 - 4:15")


# 6. Closing / next steps
s = prs.slides.add_slide(BLANK)
bg(s)
kicker(s, "Next Steps", 6)
headline(s, "แผนต่อไป (Phase 2)", "งานที่ยังค้างและจะดำเนินการต่อ เรียงตาม priority")

stat(s, Inches(0.95), Inches(2.2), "ROC", "Threshold tuning", RED)
stat(s, Inches(3.45), Inches(2.2), "Redis", "Escalation persist", TEAL)
stat(s, Inches(5.95), Inches(2.2), "BLE", "WiFi provisioning", BLUE)
stat(s, Inches(8.45), Inches(2.2), "Cloud", "Production deploy", GREEN)

rounded(s, Inches(1.05), Inches(4.25), Inches(11.25), Inches(1.1), BLACK, None)
center_text(s, "งานต่อไปคือทำให้ threshold แม่นขึ้น ระบบทนต่อ restart\nและติดตั้งได้จริงในบ้านลูกค้าโดยไม่ต้อง hardcode WiFi", Inches(1.35), Inches(4.47), Inches(10.65), Inches(0.55), 22, WHITE, True)

rounded(s, Inches(2.0), Inches(6.0), Inches(9.3), Inches(0.52), rgb("#EAF4FF"), None)
center_text(s, "Mentor feedback รอบถัดไป (ยังไม่เริ่ม): unseen test, one-class model, realistic dataset, video demo", Inches(2.2), Inches(6.14), Inches(8.9), Inches(0.2), 12.5, BLUE, True, LATIN)
footer(s, "4:15 - 5:00")


prs.save(OUT)
print(f"saved {OUT}")
