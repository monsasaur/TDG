"""Generate CS423 fall-detection presentation (15 slides, Thai, 16:9).

Usage:
    pip install python-pptx
    python generate_slides.py
    # output: fall_detection_cs423.pptx (next to this script)

Theme: Tech Minimalist
    bg=#F8F9FA  navy=#0A2540  cyan=#00B4D8  teal=#06A77D  amber=#F77F00
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------- theme ----------
OFFWHITE = RGBColor(0xF8, 0xF9, 0xFA)
NAVY = RGBColor(0x0A, 0x25, 0x40)
CYAN = RGBColor(0x00, 0xB4, 0xD8)
TEAL = RGBColor(0x06, 0xA7, 0x7D)
AMBER = RGBColor(0xF7, 0x7F, 0x00)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_CYAN = RGBColor(0xD9, 0xF1, 0xF7)
LIGHT_TEAL = RGBColor(0xD7, 0xEE, 0xE5)
LIGHT_AMBER = RGBColor(0xFD, 0xE5, 0xCC)
LIGHT_GRAY = RGBColor(0xE9, 0xEC, 0xEF)
DARK_BG = RGBColor(0x12, 0x1E, 0x2D)
CODE_GREEN = RGBColor(0x7E, 0xE0, 0xC4)

FONT_TH = "Sarabun"
FONT_MONO = "Consolas"
TOTAL = 15


# ---------- presentation ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None, line_w=0.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    return s


def set_bg(slide, color=OFFWHITE):
    bg = add_rect(slide, 0, 0, SW, SH, color)
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _run(p, text, font=FONT_TH, size=14, color=NAVY, bold=False, italic=False):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return r


def add_text(slide, x, y, w, h, text, *, font=FONT_TH, size=14, color=NAVY,
             bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    _run(p, text, font=font, size=size, color=color, bold=bold, italic=italic)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=NAVY,
                line_spacing=1.4, bold=False):
    """items: list of strings, or list of (text, color, bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            txt, c, b = item
            _run(p, txt, size=size, color=c, bold=b)
        else:
            _run(p, item, size=size, color=color, bold=bold)
    return tb


def add_header(slide, title, n):
    add_rect(slide, 0, 0, Inches(0.18), SH, CYAN)
    add_text(slide, Inches(0.55), Inches(0.32), Inches(12.5), Inches(0.65),
             title, size=26, bold=True, color=NAVY)
    add_rect(slide, Inches(0.55), Inches(0.96), Inches(2.4), Emu(12000), NAVY)
    add_text(slide, Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.3),
             f"{n} / {TOTAL} · CS423", size=10, color=GRAY,
             align=PP_ALIGN.RIGHT)


def add_callout(slide, x, y, w, h, text, *, fill=CYAN, color=WHITE,
                size=14, bold=True, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3):
    box = add_rect(slide, x, y, w, h, fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        _run(p, line, size=size, color=color, bold=bold)
    return box


def add_table(slide, x, y, w, h, headers, rows, *,
              col_widths=None, font_size=12,
              header_fill=NAVY, header_color=WHITE,
              highlight_row=None, highlight_col=None,
              highlight_fill=LIGHT_CYAN, cell_aligns=None):
    cols = len(headers)
    nrows = len(rows) + 1
    shape = slide.shapes.add_table(nrows, cols, x, y, w, h)
    t = shape.table
    if col_widths:
        total = sum(col_widths)
        for j, cw in enumerate(col_widths):
            t.columns[j].width = int(w * cw / total)

    def style_cell(cell, txt, *, bold=False, color=NAVY, size=font_size,
                   fill=None, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill if fill is not None else WHITE
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        _run(p, str(txt), size=size, color=color, bold=bold)

    for j, htxt in enumerate(headers):
        style_cell(t.cell(0, j), htxt, bold=True, color=header_color,
                   fill=header_fill)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            fill = None
            if highlight_row is not None and i == highlight_row:
                fill = highlight_fill
            elif highlight_col is not None and j == highlight_col:
                fill = highlight_fill
            align = PP_ALIGN.LEFT
            if cell_aligns and j < len(cell_aligns):
                align = cell_aligns[j]
            style_cell(t.cell(i + 1, j), val, fill=fill, align=align)
    return t


def add_image_box(slide, x, y, w, h, label, *, size=12):
    box = add_rect(slide, x, y, w, h, LIGHT_GRAY, line=GRAY, line_w=0.75)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, f"[ {label} ]", size=size, italic=True, color=GRAY)
    return box


def add_code_block(slide, x, y, w, h, lines, *, size=12,
                   color=CODE_GREEN, fill=DARK_BG, highlights=None):
    """lines: list[str]; highlights: dict[int, RGBColor] for per-line color."""
    add_rect(slide, x, y, w, h, fill)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.35
        c = (highlights or {}).get(i, color)
        _run(p, ln, font=FONT_MONO, size=size, color=c)
    return tb


def add_card(slide, x, y, w, h, *, fill=WHITE, line=LIGHT_GRAY):
    return add_rect(slide, x, y, w, h, fill, line=line, line_w=0.75)


def add_metric_card(slide, x, y, w, h, value, label, *,
                    value_color=TEAL, value_size=48, label_size=14):
    add_card(slide, x, y, w, h)
    add_text(slide, x, y + Inches(0.35), w, Inches(1.2), value,
             size=value_size, bold=True, color=value_color,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + h - Inches(0.7), w, Inches(0.5), label,
             size=label_size, color=NAVY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


# ---------- slide 1: cover ----------
def slide_1():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_rect(s, 0, 0, Inches(0.6), SH, CYAN)
    add_text(s, Inches(1.2), Inches(1.4), Inches(11.5), Inches(1.0),
             "ระบบตรวจจับการล้มของผู้สูงอายุ",
             size=42, bold=True, color=NAVY)
    add_text(s, Inches(1.2), Inches(2.3), Inches(11.5), Inches(0.9),
             "ด้วย IoT และ WiFi RF Sensing",
             size=34, bold=True, color=NAVY)
    add_text(s, Inches(1.2), Inches(3.3), Inches(11.5), Inches(0.5),
             "Fall Detection IoT System using WiFi CSI",
             size=18, italic=True, color=CYAN)
    add_rect(s, Inches(1.2), Inches(3.95), Inches(3), Emu(15000), NAVY)
    add_text(s, Inches(1.2), Inches(4.25), Inches(6), Inches(0.35),
             "จัดทำโดย", size=13, bold=True, color=NAVY)
    add_text(s, Inches(1.2), Inches(4.62), Inches(8), Inches(0.4),
             "[ ใส่ชื่อ-นามสกุลและรหัสนักศึกษา ]",
             size=14, color=GRAY)
    add_text(s, Inches(1.2), Inches(5.5), Inches(11), Inches(0.4),
             "เสนอ: อ.ทศพร เวชศิริ", size=14, color=NAVY)
    add_text(s, Inches(1.2), Inches(5.92), Inches(11), Inches(0.4),
             "รายวิชา: CS423 Internet of Things and Applications",
             size=14, color=NAVY)
    add_text(s, Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.3),
             "1 / 15 · CS423", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ---------- slide 2: ที่มาและความสำคัญ ----------
def slide_2():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "ที่มาและความสำคัญ", 2)

    add_text(s, Inches(0.55), Inches(1.3), Inches(6), Inches(0.4),
             "สถานการณ์", size=16, bold=True, color=CYAN)
    add_bullets(s, Inches(0.7), Inches(1.75), Inches(12), Inches(1.3), [
        "•  ผู้สูงอายุ 1 ใน 3 ที่อายุเกิน 60 ปี ประสบเหตุล้มทุกปี",
        "•  การช่วยเหลือล่าช้าเกิน 3 ชั่วโมง อาจถึงแก่ชีวิต",
    ], size=15)

    add_text(s, Inches(0.55), Inches(3.05), Inches(8), Inches(0.4),
             "ข้อจำกัดของระบบที่มีอยู่", size=16, bold=True, color=CYAN)
    add_table(s, Inches(0.7), Inches(3.5), Inches(12), Inches(2.0),
              ["ระบบ", "ข้อจำกัด"],
              [
                  ["กล้องวงจรปิด", "ละเมิดความเป็นส่วนตัว"],
                  ["อุปกรณ์สวมใส่", "ลืมสวม / ต้องชาร์จ"],
                  ["เซนเซอร์แรงดันบนพื้น", "ติดตั้งยาก / ครอบคลุมจำกัด"],
              ],
              col_widths=[1, 2], font_size=13)

    add_callout(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.85),
                "▶  แนวทางของเรา: IoT แบบ non-intrusive ใช้คลื่น WiFi เป็น sensing medium",
                fill=CYAN, color=WHITE, size=15, bold=True)


# ---------- slide 3: RF Sensing ----------
def slide_3():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "RF Sensing — เทคนิค IoT ขั้นสูง", 3)

    add_text(s, Inches(0.7), Inches(1.25), Inches(12), Inches(0.9),
             "RF Sensing คือการใช้คลื่นวิทยุที่มีอยู่ในสภาพแวดล้อมเป็นเซนเซอร์ "
             "ไม่ต้องติดตั้งเซนเซอร์เพิ่ม ไม่ต้องแปะอุปกรณ์ที่ตัวคน",
             size=14, color=NAVY, line_spacing=1.4)

    add_text(s, Inches(0.55), Inches(2.25), Inches(8), Inches(0.4),
             "เปรียบเทียบ Wireless Sensing", size=15, bold=True, color=CYAN)
    add_table(s, Inches(0.7), Inches(2.7), Inches(12), Inches(1.95),
              ["เทคโนโลยี", "อุปกรณ์", "ครอบคลุม"],
              [
                  ["Bluetooth RSSI", "BLE beacon", "ใกล้"],
                  ["UWB", "Anchor 3-4 ตัว", "กลาง"],
                  ["Lidar", "Lidar sensor", "แคบ"],
                  ["WiFi CSI (เรา) ✓", "ESP32 ×2", "ทั้งห้อง"],
              ],
              col_widths=[1, 1.2, 1], font_size=13,
              highlight_row=3, highlight_fill=LIGHT_TEAL)

    add_text(s, Inches(0.55), Inches(4.85), Inches(8), Inches(0.4),
             "CSI Pattern แต่ละกิจกรรม", size=15, bold=True, color=CYAN)
    w = Inches(3.95)
    gap = Inches(0.18)
    x0 = Inches(0.7)
    for i, label in enumerate(["CSI: ยืนนิ่ง", "CSI: เดิน", "CSI: ล้ม"]):
        add_image_box(s, x0 + (w + gap) * i, Inches(5.3), w, Inches(1.6), label)


# ---------- slide 4: hardware components ----------
def slide_4():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "Hardware Components", 4)

    add_text(s, Inches(0.55), Inches(1.25), Inches(8), Inches(0.4),
             "อุปกรณ์ที่ใช้ทั้งหมด", size=15, bold=True, color=CYAN)
    add_table(s, Inches(0.7), Inches(1.7), Inches(12), Inches(2.1),
              ["รายการ", "จำนวน", "หน้าที่"],
              [
                  ["ESP32-WROOM-32", "×2", "MCU + WiFi CSI extraction"],
                  ["USB-C Cable", "×2", "จ่ายไฟ + flash firmware"],
                  ["Power Bank 5V/2A", "×2", "จ่ายไฟตอน deploy"],
                  ["Smartphone", "×1", "Hotspot + Mobile App"],
                  ["Laptop / Server", "×1", "รับ UDP + ML service"],
              ],
              col_widths=[1.4, 0.5, 2], font_size=12,
              highlight_row=0, highlight_fill=LIGHT_CYAN,
              cell_aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])

    cw = Inches(5.85)
    ch = Inches(1.85)
    gap = Inches(0.3)
    x_l = Inches(0.7)
    x_r = x_l + cw + gap
    y = Inches(4.0)

    add_card(s, x_l, y, cw, ch)
    add_image_box(s, x_l + Inches(0.15), y + Inches(0.15),
                  Inches(1.4), ch - Inches(0.3), "รูป ESP32 #1", size=10)
    add_text(s, x_l + Inches(1.7), y + Inches(0.2), cw - Inches(1.85), Inches(0.4),
             "ESP32 #1 (active_ap)", size=14, bold=True, color=NAVY)
    add_bullets(s, x_l + Inches(1.7), y + Inches(0.7), cw - Inches(1.85), Inches(1.1), [
        "•  Dual mode AP+STA",
        "•  วัด CSI + ส่ง UDP",
    ], size=12, line_spacing=1.4)

    add_card(s, x_r, y, cw, ch)
    add_image_box(s, x_r + Inches(0.15), y + Inches(0.15),
                  Inches(1.4), ch - Inches(0.3), "รูป ESP32 #2", size=10)
    add_text(s, x_r + Inches(1.7), y + Inches(0.2), cw - Inches(1.85), Inches(0.4),
             "ESP32 #2 (active_sta)", size=14, bold=True, color=NAVY)
    add_bullets(s, x_r + Inches(1.7), y + Inches(0.7), cw - Inches(1.85), Inches(1.1), [
        "•  STA only",
        "•  ส่ง WiFi packet ให้ตัววัด",
    ], size=12, line_spacing=1.4)

    add_callout(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.7),
                "💰  ต้นทุนรวมต่อชุด ~ 600 บาท",
                fill=AMBER, color=WHITE, size=15, bold=True)


# ---------- slide 5: ทำไม ESP32 ----------
def slide_5():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "ทำไม ESP32? — Hardware Specification", 5)

    add_text(s, Inches(0.55), Inches(1.2), Inches(8), Inches(0.4),
             "เปรียบเทียบกับ MCU ตัวอื่น", size=14, bold=True, color=CYAN)
    add_table(s, Inches(0.7), Inches(1.6), Inches(12), Inches(2.4),
              ["คุณสมบัติ", "ESP32", "ESP8266", "Arduino", "RPi"],
              [
                  ["CSI extraction", "✅", "❌", "❌", "❌"],
                  ["Dual-core 240MHz", "✅", "❌", "❌", "✅"],
                  ["WiFi onboard", "✅", "✅", "❌", "✅"],
                  ["Dual-mode AP+STA", "✅", "จำกัด", "❌", "✅"],
                  ["ราคา (บาท)", "250", "150", "300", "1,500"],
                  ["Power (mA)", "80", "70", "50", "600"],
              ],
              col_widths=[2, 1, 1, 1, 1], font_size=11,
              highlight_col=1, highlight_fill=LIGHT_TEAL,
              cell_aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER,
                           PP_ALIGN.CENTER, PP_ALIGN.CENTER])

    add_text(s, Inches(0.55), Inches(4.15), Inches(10), Inches(0.4),
             "ESP32-WROOM-32 Specification", size=14, bold=True, color=CYAN)
    add_code_block(s, Inches(0.7), Inches(4.55), Inches(12), Inches(2.1), [
        "Chip      : ESP32-D0WD-V3",
        "CPU       : Xtensa LX6 Dual-Core @ 240 MHz",
        "Memory    : 520 KB SRAM + 448 KB ROM | Flash 4 MB",
        "WiFi      : 802.11 b/g/n (2.4 GHz, 150 Mbps)",
        "Bluetooth : v4.2 BR/EDR + BLE",
        "GPIO      : 34 ขา | Antenna: PCB onboard",
        "Voltage   : 3.3V (USB 5V → AMS1117) | Current 80-260 mA TX",
    ], size=11)

    add_text(s, Inches(0.7), Inches(6.75), Inches(12), Inches(0.35),
             "▶ จุดตัดสินใจ: ESP32 เป็น MCU ตัวเดียวที่ดึง CSI ได้ในราคาเข้าถึงได้",
             size=14, bold=True, color=CYAN)


# ---------- slide 6: การประกอบและการเชื่อมต่อ ----------
def slide_6():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "การประกอบและการเชื่อมต่อ", 6)

    add_text(s, Inches(0.55), Inches(1.2), Inches(8), Inches(0.4),
             "Wiring & Network Diagram", size=14, bold=True, color=CYAN)
    add_image_box(s, Inches(0.7), Inches(1.6), Inches(12), Inches(2.0),
                  "Network Topology — ESP32#2 → packet → ESP32#1 → "
                  "Hotspot \"View\" → UDP:5500 → Mac    "
                  "|    Power Bank ↓ ESP32 ทั้ง 2")

    add_text(s, Inches(0.55), Inches(3.75), Inches(8), Inches(0.4),
             "ขั้นตอนการประกอบ", size=14, bold=True, color=CYAN)
    add_bullets(s, Inches(0.7), Inches(4.15), Inches(12), Inches(2.0), [
        "①  เสียบ USB-C เข้า ESP32 ทั้ง 2 ตัว → จ่ายไฟ Power Bank",
        "②  Flash firmware: active_ap (ESP32 #1), active_sta (ESP32 #2)",
        "③  ตั้ง Hotspot ชื่อ \"View\" บนมือถือ → ESP32 #1 auto-connect",
        "④  เปิด server บน Laptop → รับ UDP port 5500",
        "⑤  เริ่มทำงาน — ไม่มีสายไฟเชื่อมระหว่างอุปกรณ์",
    ], size=13, line_spacing=1.45)

    add_callout(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
                "💡  จุดเด่น: Wireless ทุก layer · ติดตั้งใน 5 นาที",
                fill=AMBER, color=WHITE, size=14, bold=True)


# ---------- slide 7: live demo hardware ----------
def slide_7():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "🔴 LIVE DEMO — Hardware", 7)

    add_image_box(s, Inches(0.7), Inches(1.25), Inches(12), Inches(3.4),
                  "Screenshot: Terminal UDP stream + กราฟ CSI real-time")

    add_text(s, Inches(0.55), Inches(4.8), Inches(8), Inches(0.4),
             "สิ่งที่จะสาธิต", size=14, bold=True, color=CYAN)
    add_bullets(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.4), [
        "①  ESP32 boot และเชื่อม WiFi อัตโนมัติ",
        "②  UDP packet streaming 100 packet/วินาที",
        "③  กราฟ CSI 52 subcarriers แสดง real-time",
        "④  สัญญาณเปลี่ยนเมื่อมีการเคลื่อนไหว (ยืน → เดิน → ล้ม)",
    ], size=13, line_spacing=1.35)

    add_code_block(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.45), [
        "Protocol stack: WiFi 802.11n → UDP → Server (Node.js)",
    ], size=12)


# ---------- slide 8: platform overview ----------
def slide_8():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "Platform Overview", 8)

    add_text(s, Inches(0.55), Inches(1.2), Inches(8), Inches(0.4),
             "End-to-End Data Flow", size=14, bold=True, color=CYAN)
    nodes = [
        ("ESP32", NAVY, WHITE),
        ("API\n(Express)", NAVY, WHITE),
        ("ML Model\n(LSTM)", NAVY, WHITE),
        ("🤖 AI Agent", AMBER, WHITE),
        ("DB\n(Supabase)", NAVY, WHITE),
        ("Mobile App", NAVY, WHITE),
        ("Twilio", NAVY, WHITE),
    ]
    n = len(nodes)
    bw = Inches(1.45)
    bh = Inches(0.95)
    arrow_w = Inches(0.25)
    total_w = bw * n + arrow_w * (n - 1)
    x0 = (SW - total_w) // 2
    y_box = Inches(1.7)
    for i, (label, fill, color) in enumerate(nodes):
        x = x0 + (bw + arrow_w) * i
        add_rect(s, x, y_box, bw, bh, fill)
        add_text(s, x, y_box, bw, bh, label, size=11, bold=True,
                 color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 line_spacing=1.1)
        if i < n - 1:
            ax = x + bw
            add_text(s, ax, y_box, arrow_w, bh, "→", size=18, bold=True,
                     color=GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(s, Inches(0.55), Inches(3.15), Inches(8), Inches(0.4),
             "ฟีเจอร์หลัก", size=14, bold=True, color=CYAN)
    add_bullets(s, Inches(0.7), Inches(3.55), Inches(7.5), Inches(2.7), [
        "✓  Real-time event streaming (Supabase realtime subscription)",
        "✓  🤖 Smart Alert Agent — context-aware decision making",
        "✓  Mobile push notification + acknowledge system",
        "✓  Auto emergency call ผ่าน Twilio (เสียงจาก AI gen)",
    ], size=13, line_spacing=1.55)

    add_image_box(s, Inches(8.5), Inches(3.4), Inches(4.2), Inches(3.4),
                  "Screenshot: Mobile app 2 หน้า")


# ---------- slide 9: AI Agent ----------
def slide_9():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "🤖 AI Agent — Smart Alert Orchestrator", 9)

    add_text(s, Inches(0.55), Inches(1.2), Inches(6), Inches(0.4),
             "Flow การทำงาน", size=14, bold=True, color=CYAN)

    steps = [
        ("1.  LSTM detect fall (97.9%)", NAVY, False),
        ("2.  🤖 Claude Haiku วิเคราะห์ context", NAVY, True),
        ("       •  เวลาปัจจุบัน", GRAY, False),
        ("       •  Event history", GRAY, False),
        ("       •  ผู้ดูแล / ความเร่งด่วน", GRAY, False),
        ("3.  ตัดสินใจ: โทรใคร / ข้อความเสียง / timing", NAVY, False),
        ("4.  Twilio โทรพร้อม personalized voice", NAVY, False),
    ]
    add_bullets(s, Inches(0.7), Inches(1.65), Inches(6), Inches(4.8),
                steps, size=13, line_spacing=1.55)

    add_text(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
             "ทำไมไม่ใช้ if-else?", size=14, bold=True, color=CYAN)
    add_table(s, Inches(7.0), Inches(1.65), Inches(5.9), Inches(2.85),
              ["If-else (ระบบเดิม)", "AI Agent (ของเรา)"],
              [
                  ["Rule ตายตัว", "Context-aware"],
                  ["ข้อความเหมือนเดิมทุกเคส", "Personalized"],
                  ["เพิ่ม rule = แก้โค้ด", "ปรับด้วย prompt"],
              ],
              col_widths=[1, 1], font_size=12,
              highlight_col=1, highlight_fill=LIGHT_TEAL)

    add_callout(s, Inches(7.0), Inches(4.85), Inches(5.9), Inches(1.7),
                "ผลลัพธ์: ระบบยืดหยุ่น ปรับพฤติกรรมได้โดยไม่ต้อง redeploy "
                "และตอบสนองได้เหมาะกับสถานการณ์จริง",
                fill=CYAN, color=WHITE, size=13, bold=False, align=PP_ALIGN.LEFT)


# ---------- slide 10: architecture stack ----------
def slide_10():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "Architecture Stack", 10)

    add_image_box(s, Inches(0.7), Inches(1.25), Inches(12), Inches(2.5),
                  "Architecture Diagram — ESP32 ×2 → Express API → "
                  "FastAPI+LSTM / Claude Haiku Agent / Supabase / Twilio "
                  "→ React Native App")

    rows = [
        ("Layer", "Technology", None),
        ("Device", "ESP32 ×2 (ESP-IDF, dual-mode AP+STA)", None),
        ("Communication", "WiFi 802.11n, UDP, HTTPS, REST", None),
        ("Backend / Cloud", "Express.js, FastAPI, Render", None),
        ("Database", "Supabase (PostgreSQL + Realtime)", None),
        ("Frontend", "React Native + Expo", None),
        ("AI Agent (Runtime)", "Claude Haiku — Smart Alert Orchestrator",
         LIGHT_AMBER),
        ("AI Agent (Dev Tool)", "Claude Code — พัฒนา codebase", LIGHT_AMBER),
    ]
    # build with custom highlight per row
    nrows = len(rows)
    y_tbl = Inches(3.95)
    h_tbl = Inches(2.95)
    w_tbl = Inches(12)
    shape = s.shapes.add_table(nrows, 2, Inches(0.7), y_tbl, w_tbl, h_tbl)
    t = shape.table
    t.columns[0].width = int(w_tbl * 1.0 / 3.6)
    t.columns[1].width = int(w_tbl * 2.6 / 3.6)
    for i, (a, b, fill) in enumerate(rows):
        is_header = i == 0
        for j, val in enumerate((a, b)):
            cell = t.cell(i, j)
            cell.fill.solid()
            if is_header:
                cell.fill.fore_color.rgb = NAVY
            elif fill is not None:
                cell.fill.fore_color.rgb = fill
            else:
                cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            _run(p, val, size=12,
                 color=WHITE if is_header else NAVY,
                 bold=is_header)


# ---------- slide 11: live demo platform ----------
def slide_11():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "🔴 LIVE DEMO — Platform End-to-End", 11)

    add_image_box(s, Inches(0.7), Inches(1.3), Inches(5.0), Inches(4.6),
                  "Screenshot: Mobile app\nหน้า notification")

    add_bullets(s, Inches(6.0), Inches(1.3), Inches(6.7), Inches(4.6), [
        "①  จำลองล้มหน้า ESP32",
        "②  UDP → API (latency < 100ms)",
        "③  LSTM classify → fall (98% confidence)",
        "④  🤖 AI Agent วิเคราะห์ + ตัดสินใจ",
        "⑤  Mobile app ขึ้น notification ทันที",
        "⑥  ไม่กดยืนยัน 30 วิ → Twilio โทรจริง พร้อมเสียงที่ AI gen",
    ], size=14, line_spacing=1.65)

    add_text(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.7),
             "Latency รวม end-to-end: < 1 วินาที",
             size=32, bold=True, color=TEAL,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------- slide 12: ผลการทำงาน ----------
def slide_12():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "ผลการทำงานของระบบ", 12)

    metrics_top = [("< 1 s", "Latency"), ("100/s", "Throughput"),
                   ("24h", "Uptime")]
    metrics_bot = [("97.9%", "Accuracy"), ("98.5%", "Fall Recall"),
                   ("< 800ms", "AI Agent Response")]
    cw = Inches(3.95)
    ch = Inches(2.3)
    gap = Inches(0.18)
    x0 = Inches(0.55)
    y_top = Inches(1.4)
    y_bot = y_top + ch + gap
    for i, (v, lbl) in enumerate(metrics_top):
        add_metric_card(s, x0 + (cw + gap) * i, y_top, cw, ch, v, lbl)
    for i, (v, lbl) in enumerate(metrics_bot):
        add_metric_card(s, x0 + (cw + gap) * i, y_bot, cw, ch, v, lbl)

    add_text(s, Inches(0.55), Inches(6.55), Inches(12.2), Inches(0.4),
             "Dataset: 1,260 windows · เก็บจาก ESP32 จริงทั้งหมด",
             size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# ---------- slide 13: ทำไม hardware น้อย ----------
def slide_13():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "ทำไม Hardware น้อย?", 13)

    add_text(s, Inches(0.55), Inches(1.2), Inches(8), Inches(0.4),
             "RF Sensing vs Sensor Array", size=14, bold=True, color=CYAN)

    add_table(s, Inches(0.7), Inches(1.65), Inches(12), Inches(3.0),
              ["", "Sensor Array (ทั่วไป)", "RF Sensing (เรา) ✓"],
              [
                  ["จำนวนอุปกรณ์", "5-10 ตัว", "2 ตัว"],
                  ["ติดตั้ง", "หลายจุดในห้อง", "มุมเดียว"],
                  ["Protocol stack", "I2C/SPI → MCU → WiFi", "WiFi ทุก layer"],
                  ["ครอบคลุม", "จำกัดตามจุดที่วาง", "ทั้งห้อง"],
                  ["Complexity", "Wiring + Calibration",
                   "Network + Real-time + AI"],
              ],
              col_widths=[1.1, 1.4, 1.4], font_size=12,
              highlight_col=2, highlight_fill=LIGHT_TEAL)

    add_callout(s, Inches(0.7), Inches(5.0), Inches(12), Inches(1.85),
                "💡  Hardware น้อย ≠ ง่าย\n"
                "Complexity ย้ายไปอยู่ที่ Network Protocol Stack + AI Integration",
                fill=AMBER, color=WHITE, size=18, bold=True,
                align=PP_ALIGN.CENTER, line_spacing=1.5)


# ---------- slide 14: tech stack + roadmap ----------
def slide_14():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_header(s, "Tech Stack + Roadmap", 14)

    code_lines = [
        "Firmware       ESP-IDF + C (dual-mode AP+STA, CSI extraction)",
        "Networking     WiFi 802.11n + UDP streaming",
        "Database       Supabase (PostgreSQL + Realtime)",
        "Frontend       React Native + Expo",
        "Telecom        Twilio Voice API",
        "AI Runtime     Claude Haiku — Smart Alert Agent",
        "AI Dev Tool    Claude Code — boilerplate + debugging",
    ]
    add_code_block(s, Inches(0.7), Inches(1.25), Inches(12), Inches(2.4),
                   code_lines, size=11,
                   highlights={5: AMBER, 6: AMBER})

    add_text(s, Inches(0.55), Inches(3.85), Inches(6), Inches(0.4),
             "ปัญหา + วิธีแก้", size=14, bold=True, color=CYAN)
    add_bullets(s, Inches(0.7), Inches(4.25), Inches(6.0), Inches(2.6), [
        "•  UDP packet loss → sliding window + stride",
        "•  IP เปลี่ยนตอน hotspot reconnect → hardcode IP firmware",
        "•  Mobile cold-start lag → Supabase realtime subscription",
        "•  AI Agent latency → ใช้ prompt caching ลด 90%",
    ], size=12, line_spacing=1.55)

    add_text(s, Inches(7.0), Inches(3.85), Inches(6), Inches(0.4),
             "Roadmap", size=14, bold=True, color=CYAN)
    add_bullets(s, Inches(7.15), Inches(4.25), Inches(5.7), Inches(2.6), [
        "🔧  BLE Provisioning — ลูกค้าตั้ง WiFi เอง",
        "🏠  Multi-room ESP32 mesh",
        "🎯  On-device inference บน ESP32-S3",
        "🤖  Voice Conversation Agent",
        "💼  Business: บริการบ้านพักคนชรา ~600฿/บ้าน",
    ], size=12, line_spacing=1.55)


# ---------- slide 15: Q&A ----------
def slide_15():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_rect(s, 0, 0, Inches(0.6), SH, CYAN)

    add_text(s, Inches(0.55), Inches(2.1), Inches(12.3), Inches(1.2),
             "ขอบคุณครับ", size=64, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    line_w = Inches(1.5)
    add_rect(s, (SW - line_w) // 2, Inches(3.35), line_w, Emu(15000), NAVY)
    add_text(s, Inches(0.55), Inches(3.5), Inches(12.3), Inches(0.7),
             "Questions?", size=32, italic=True, color=CYAN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    qr_w = Inches(1.7)
    qr_y = Inches(4.6)
    qr_left = (SW // 2) - Inches(2.0)
    qr_right = (SW // 2) + Inches(0.3)
    add_image_box(s, qr_left, qr_y, qr_w, qr_w, "QR\nGitHub Repo", size=10)
    add_image_box(s, qr_right, qr_y, qr_w, qr_w, "QR\nDemo Video", size=10)

    add_text(s, Inches(0.55), Inches(6.6), Inches(12.3), Inches(0.4),
             "[ ใส่ email ติดต่อทีม ]", size=12, color=GRAY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.55), Inches(7.08), Inches(12.2), Inches(0.3),
             "15 / 15 · CS423", size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ---------- build ----------
def main():
    for fn in (slide_1, slide_2, slide_3, slide_4, slide_5,
               slide_6, slide_7, slide_8, slide_9, slide_10,
               slide_11, slide_12, slide_13, slide_14, slide_15):
        fn()
    out = Path(__file__).resolve().parent / "fall_detection_cs423.pptx"
    prs.save(out)
    print(f"saved: {out}")


if __name__ == "__main__":
    main()
